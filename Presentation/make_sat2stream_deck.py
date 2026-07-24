#!/usr/bin/env python3
"""
Sat2Stream_progress.pptx — job-talk style rebuild.
Design language: white slides, one accent color (#2a78d6), kicker + large
title + thin accent rule, figures dominate, one claim per slide.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

HERE = os.path.dirname(os.path.abspath(__file__))
A = os.path.join(HERE, 'sat2stream_assets')
OUT = os.path.join(HERE, 'Sat2Stream_progress.pptx')

SW, SH = Inches(13.333), Inches(7.5)
FONT = 'Liberation Sans'
ACCENT = RGBColor(0x2A, 0x78, 0xD6)
ORANGE = RGBColor(0xEB, 0x68, 0x34)
INK = RGBColor(0x0B, 0x0B, 0x0B)
INK2 = RGBColor(0x52, 0x51, 0x4E)
FAINT = RGBColor(0x9A, 0x99, 0x94)
LINE = RGBColor(0xD8, 0xD7, 0xD3)

prs = Presentation()
prs.slide_width, prs.slide_height = SW, SH
BLANK = prs.slide_layouts[6]


def F(run, size, bold=False, color=INK, italic=False):
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color


def text(s, x, y, w, h, lines, size=15, bold=False, color=INK,
         align=PP_ALIGN.LEFT, spacing=1.18, anchor=MSO_ANCHOR.TOP):
    box = s.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    if isinstance(lines, str):
        lines = [lines]
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        if isinstance(ln, tuple):
            t, sz, b, c, *rest = ln
            it = rest[0] if rest else False
        else:
            t, sz, b, c, it = ln, size, bold, color, False
        r = p.add_run()
        r.text = t
        F(r, sz, b, c, it)
    return box


def rule(s, x, y, w=Inches(1.05), color=ACCENT, h=Pt(3.2)):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Inches(0.045))
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    return sh


def header(s, kicker, title, tsize=27):
    text(s, Inches(0.62), Inches(0.34), Inches(11), Inches(0.35), kicker.upper(),
         size=13, bold=True, color=FAINT)
    text(s, Inches(0.6), Inches(0.66), Inches(12.2), Inches(0.75), title,
         size=tsize, bold=True, color=INK)
    rule(s, Inches(0.64), Inches(1.38))


def footer(s, n):
    text(s, Inches(0.62), SH - Inches(0.42), Inches(4), Inches(0.3),
         'Sat2Stream — progress', size=10, color=FAINT)
    text(s, SW - Inches(0.9), SH - Inches(0.42), Inches(0.5), Inches(0.3),
         str(n), size=11, color=FAINT, align=PP_ALIGN.RIGHT)


def panel(s, x, y, w, h, edge=LINE, weight=1.0):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sh.adjustments[0] = 0.06
    sh.fill.solid()
    sh.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    sh.line.color.rgb = edge
    sh.line.width = Pt(weight)
    sh.shadow.inherit = False
    return sh


def arrow(s, x1, y1, x2, y2, color=INK2, weight=1.8, dash=False):
    ln = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    ln.line.color.rgb = color
    ln.line.width = Pt(weight)
    el = ln.line._get_or_add_ln()
    el.append(el.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'}))
    if dash:
        d = el.makeelement(qn('a:prstDash'), {'val': 'dash'})
        el.insert(0, d)
    return ln


def slide():
    return prs.slides.add_slide(BLANK)


# ================= 1 · TITLE ==============================================
s = slide()
text(s, Inches(1.0), Inches(2.05), Inches(11.3), Inches(1.1), 'Sat2Stream',
     size=52, bold=True, align=PP_ALIGN.CENTER)
rule(s, Inches(6.17), Inches(3.12), w=Inches(1.0))
text(s, Inches(1.0), Inches(3.35), Inches(11.3), Inches(0.6),
     'Multimodal Flow Matching for Zero-Shot Streamflow Reconstruction',
     size=21, color=INK2, align=PP_ALIGN.CENTER)
text(s, Inches(1.0), Inches(4.25), Inches(11.3), Inches(0.5),
     [('What satellites see, flow matching restores.', 16, False, INK2, True)],
     align=PP_ALIGN.CENTER)
text(s, Inches(1.0), Inches(5.6), Inches(11.3), Inches(0.9), [
    ('Progress report · group meeting · July 2026', 14, False, FAINT),
    ('builds on ZeroDiff (ICML 2026 submission)', 13, False, FAINT),
], align=PP_ALIGN.CENTER, spacing=1.35)

# ================= 2 · PROBLEM ============================================
s = slide()
header(s, 'the problem', 'Reconstruct streamflow where it was never measured')
py, ph = Inches(2.1), Inches(2.5)
panel(s, Inches(1.0), py, Inches(4.6), ph)
text(s, Inches(1.35), py + Inches(0.3), Inches(4.0), Inches(2.0), [
    ('Gauged basins', 18, True, INK),
    ('', 6, False, INK),
    ('X  drivers everywhere      ✓', 15, False, INK2),
    ('Y  streamflow observed   ✓', 15, False, INK2),
    ('', 6, False, INK),
    ('→ training', 15, True, ACCENT),
], spacing=1.25)
panel(s, Inches(7.7), py, Inches(4.6), ph)
text(s, Inches(8.05), py + Inches(0.3), Inches(4.0), Inches(2.0), [
    ('Ungauged basins', 18, True, INK),
    ('', 6, False, INK),
    ('X  drivers everywhere      ✓', 15, False, INK2),
    ('Y  never measured           —', 15, True, ORANGE),
    ('', 6, False, INK),
    ('→ reconstruct the whole series', 15, True, ACCENT),
], spacing=1.25)
arrow(s, Inches(5.68), py + Inches(1.25), Inches(7.62), py + Inches(1.25),
      color=ACCENT, weight=2.6)
text(s, Inches(5.68), py + Inches(0.75), Inches(1.95), Inches(0.4), 'zero-shot',
     size=13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
text(s, Inches(1.0), Inches(5.35), Inches(11.3), Inches(1.0), [
    ('Not forecasting (needs same-site history).  Not imputation (needs partial observations).', 15, False, INK2),
    ('The model must infer dynamics it has never seen at this location — with unknown target statistics.', 15, False, INK2),
], spacing=1.4)
footer(s, 2)

# ================= 3 · ZERODIFF ===========================================
s = slide()
header(s, 'starting point', 'ZeroDiff: calibrate a structured prior, don’t generate from noise')
s.shapes.add_picture(os.path.join(A, 'framework-1.png'),
                     Inches(0.62), Inches(1.75), width=Inches(7.5))
tx = Inches(8.5)
text(s, tx, Inches(1.95), Inches(4.3), Inches(4.2), [
    ('1 · Informed prior', 17, True, INK),
    ('LSTM maps drivers X to a structured estimate ŷ; a VAE infers each basin’s target statistics from X alone.', 14, False, INK2),
    ('', 10, False, INK),
    ('2 · Generative calibration', 17, True, INK),
    ('Diffusion / stochastic interpolant learns the transport ŷ → y, conditioned on X.', 14, False, INK2),
    ('', 10, False, INK),
    ('3 · Zero-shot transfer', 17, True, INK),
    ('Trained on gauged basins; applied unchanged at basins with no labels.', 14, False, INK2),
], spacing=1.22)
footer(s, 3)

# ================= 4 · GAP ================================================
s = slide()
header(s, 'the gap', 'The prior cannot see stored water')
s.shapes.add_picture(os.path.join(A, 'concept-1.png'),
                     Inches(0.75), Inches(1.85), width=Inches(4.7))
text(s, Inches(0.78), Inches(6.15), Inches(4.7), Inches(0.5),
     'exogenous-only prior: peaks over-smoothed (ZeroDiff, Fig. 1)',
     size=11.5, color=FAINT)
text(s, Inches(6.1), Inches(2.15), Inches(6.6), Inches(3.6), [
    ('Ground drivers describe the weather —', 21, True, INK),
    ('not the watershed’s stored water.', 21, True, INK),
    ('', 12, False, INK),
    ('Runoff response depends on antecedent storage:', 15, False, INK2),
    ('how wet the basin already is when the rain arrives.', 15, False, INK2),
    ('', 8, False, INK),
    ('X contains no such state — the prior can only infer it,', 15, False, INK2),
    ('and inferred state over-smooths and misses extremes.', 15, False, INK2),
], spacing=1.3)
footer(s, 4)

# ================= 5 · IDEA (hero) ========================================
s = slide()
header(s, 'the idea', 'Satellites observe every watershed, every day')
s.shapes.add_picture(os.path.join(A, 'fig_cutout.png'),
                     Inches(7.15), Inches(1.95), width=Inches(5.5))
text(s, Inches(7.2), Inches(6.35), Inches(5.4), Inches(0.6),
     'Potomac River (01646500) · one day of SMAP surface soil moisture · out-of-basin pixels blank',
     size=11.5, color=FAINT)
text(s, Inches(0.75), Inches(2.3), Inches(6.1), Inches(3.6), [
    ('Turn satellites into', 26, True, INK),
    ('virtual gauges.', 26, True, ACCENT),
    ('', 12, False, INK),
    ('ZeroDiff had to infer the basin state from weather.', 15.5, False, INK2),
    ('SMAP L4 soil moisture lets us observe it —', 15.5, False, INK2),
    ('a key state variable, at 9 km, every day.', 15.5, False, INK2),
    ('', 8, False, INK),
    ('Crucially: the modality never disappears at', 15.5, False, INK2),
    ('test basins. Ungauged ≠ unobserved.', 15.5, True, INK),
], spacing=1.28)
footer(s, 5)

# ================= 6 · DATA ===============================================
s = slide()
header(s, 'data', '618 basins × a decade of daily satellite fields')
s.shapes.add_picture(os.path.join(A, 'fig_map.png'),
                     Inches(0.62), Inches(1.95), width=Inches(7.7))
text(s, Inches(0.75), Inches(6.45), Inches(7.4), Inches(0.5),
     '618 CAMELSH basins · all 18 CONUS hydrologic regions',
     size=12, color=FAINT)
sx = Inches(8.85)
stats = [('618', 'basins paired: hydromet ⊕ SMAP ⊕ attributes'),
         ('3,928', 'days (2015–2025), zero missing pixels'),
         ('18 px', 'median watershed cutout (max 873)'),
         ('9,012', 'basins available for pretraining')]
yy = Inches(2.0)
for big, small in stats:
    text(s, sx, yy, Inches(3.9), Inches(0.55), big, size=27, bold=True, color=ACCENT)
    text(s, sx, yy + Inches(0.52), Inches(3.9), Inches(0.55), small, size=13, color=INK2)
    yy += Inches(1.14)
footer(s, 6)

# ================= 7 · METHOD =============================================
s = slide()
header(s, 'method', 'Sat2Stream = ZeroDiff + a satellite branch')
by, bh = Inches(3.35), Inches(1.25)
boxes = [
    (Inches(0.8), 'Drivers  X', 'forcing + static (51)', LINE),
    (Inches(4.35), 'Stage 1 · LSTM', 'informed prior  ŷ', LINE),
    (Inches(7.9), 'Stage 2 · FM / SI', 'transport  ŷ → y', LINE),
    (Inches(11.05), 'ŷ  →  y', 'reconstructed series', LINE),
]
bw = Inches(2.7)
for x, h1, h2, edge in boxes:
    panel(s, x, by, bw if x != Inches(11.05) else Inches(1.9), bh, edge=edge)
    text(s, x + Inches(0.12), by + Inches(0.18),
         (bw if x != Inches(11.05) else Inches(1.9)) - Inches(0.24), Inches(0.9), [
             (h1, 15.5, True, INK), (h2, 12.5, False, INK2)], spacing=1.15)
arrow(s, Inches(3.55), by + Inches(0.62), Inches(4.30), by + Inches(0.62))
arrow(s, Inches(7.10), by + Inches(0.62), Inches(7.85), by + Inches(0.62))
arrow(s, Inches(10.65), by + Inches(0.62), Inches(11.00), by + Inches(0.62))
# satellite branch (accent)
pnl = panel(s, Inches(2.6), Inches(1.75), Inches(3.4), Inches(1.15), edge=ACCENT, weight=1.8)
text(s, Inches(2.74), Inches(1.9), Inches(3.15), Inches(0.9), [
    ('Satellite branch', 15.5, True, ACCENT),
    ('pixel-set encoder (DeepSets, ragged)', 12.5, False, INK2),
    ('→ one embedding per basin-day', 12.5, False, INK2)], spacing=1.15)
arrow(s, Inches(4.9), Inches(2.92), Inches(5.4), by, color=ACCENT, weight=2.2)
arrow(s, Inches(6.0), Inches(2.35), Inches(8.9), by, color=ACCENT, weight=2.2)
text(s, Inches(7.0), Inches(2.38), Inches(4.2), Inches(0.35),
     'the calibrator sees the same embedding', size=12, bold=True, color=ACCENT)
text(s, Inches(0.8), Inches(5.25), Inches(11.7), Inches(1.2), [
    ('Design for transfer:  pixels as tokens (permutation-invariant, no padding) · no attention · no contrastive alignment.', 14.5, False, INK2),
    ('The calibrator is conditioned on everything the prior consumed — X, ŷ, and the satellite embedding.', 14.5, True, INK),
], spacing=1.4)
footer(s, 7)

# ================= 8 · RESULTS ============================================
s = slide()
header(s, 'results — cross-validation fold 3 of 22 · 28 held-out basins', 'Ablation: the satellite and trust features help only in combination')
s.shapes.add_picture(os.path.join(A, 'fig_ablation.png'),
                     Inches(1.85), Inches(1.7), width=Inches(9.6))
text(s, Inches(1.9), Inches(6.62), Inches(10.0), Inches(0.7), [
    ('identical folds, seeds and time window (daily, 2015-04 → 2024-12); NSE on fully held-out basins', 12.5, False, FAINT),
    ('in-distribution, prior-only ≈ full conditioning — the difference appears under regional shift (slide 12)', 12.5, False, FAINT),
], spacing=1.25)
footer(s, 8)

# ================= 9 · ERROR ANALYSIS =====================================
s = slide()
header(s, 'error analysis', 'Three basins crashed — and alignment was not the cause')
text(s, Inches(0.8), Inches(1.9), Inches(5.6), Inches(3.9), [
    ('Ruled out, per basin:', 16, True, INK),
    ('', 6, False, INK),
    ('✓  cutout area = official drainage area', 15, False, INK2),
    ('✓  gauge inside the cutout box', 15, False, INK2),
    ('✓  sm–Q correlation ≈ 0.8, same as healthy basins', 15, False, INK2),
    ('', 10, False, INK),
    ('A misalignment bug would hurt every basin.', 15, False, INK2),
    ('This failure is specific to one geology: karst.', 15, True, INK),
], spacing=1.3)
text(s, Inches(7.1), Inches(1.8), Inches(5.4), Inches(1.3), '94%',
     size=60, bold=True, color=ORANGE)
text(s, Inches(7.15), Inches(3.0), Inches(5.4), Inches(0.5),
     'karst area at the worst crash (Antietam Creek, 01619500)',
     size=14, color=INK2)
text(s, Inches(7.1), Inches(3.8), Inches(5.5), Inches(2.4), [
    ('SMAP senses the top ~5 cm of soil.', 15.5, False, INK2),
    ('Karst stores water in bedrock aquifers,', 15.5, False, INK2),
    ('invisible from orbit.', 15.5, False, INK2),
    ('', 8, False, INK),
    ('The satellite reports “dry” while the', 15.5, True, INK),
    ('aquifer is full — the prior believes it.', 15.5, True, INK),
], spacing=1.28)
footer(s, 9)

# ================= 10 · FIX ===============================================
s = slide()
header(s, 'the fix', 'Teach the model when not to trust the sky')
s.shapes.add_picture(os.path.join(A, 'fig_rescue.png'),
                     Inches(0.7), Inches(1.85), width=Inches(6.9))
text(s, Inches(8.0), Inches(2.1), Inches(4.7), Inches(4.0), [
    ('Three static “trust” attributes:', 16, True, INK),
    ('', 6, False, INK),
    ('karst %  — storage hidden from orbit', 14.5, False, INK2),
    ('regulation %  — reservoirs decouple Q', 14.5, False, INK2),
    ('lake %  — buffering smooths response', 14.5, False, INK2),
    ('', 12, False, INK),
    ('≈ 70% of the damage recovered;', 15.5, True, INK),
    ('catastrophic failures eliminated.', 15.5, True, INK),
    ('', 8, False, INK),
    ('Residual gap vs baseline motivates', 14, False, INK2),
    ('gating and modality dropout (next).', 14, False, INK2),
], spacing=1.28)
footer(s, 10)

# ================= 11 · INTERACTION (2×2) =================================
s = slide()
header(s, 'why it works', 'The value is in the interaction')
gx, gy = Inches(2.4), Inches(2.15)
cw, ch = Inches(3.3), Inches(1.5)
for cx, htxt in [(0, '−  satellite'), (1, '+  satellite')]:
    text(s, gx + cx * (cw + Inches(0.35)), gy - Inches(0.5), cw, Inches(0.4),
         htxt, size=14.5, bold=True, color=INK2, align=PP_ALIGN.CENTER)
for cy, htxt in [(0, '−  trust'), (1, '+  trust')]:
    text(s, gx - Inches(1.45), gy + cy * (ch + Inches(0.35)) + Inches(0.5),
         Inches(1.3), Inches(0.4), htxt, size=14.5, bold=True, color=INK2,
         align=PP_ALIGN.RIGHT)
cells = [
    (0, 0, '0.530', False, None),
    (1, 0, '0.553', False, '(3 basins crash)'),
    (0, 1, '0.472', False, None),
    (1, 1, '0.576', True, None),
]
for cx, cy, val, win, note in cells:
    x = gx + cx * (cw + Inches(0.35))
    y = gy + cy * (ch + Inches(0.35))
    panel(s, x, y, cw, ch, edge=ACCENT if win else LINE,
          weight=2.2 if win else 1.0)
    text(s, x, y + Inches(0.38), cw, Inches(0.7), val, size=27, bold=True,
         color=ACCENT if win else INK, align=PP_ALIGN.CENTER)
    if note:
        text(s, x, y + Inches(1.02), cw, Inches(0.4), note, size=12,
             color=ORANGE, align=PP_ALIGN.CENTER)
text(s, gx - Inches(1.45), gy - Inches(0.5), Inches(1.3), Inches(0.4),
     'median NSE', size=12, color=FAINT, align=PP_ALIGN.RIGHT)
text(s, Inches(2.4), Inches(5.65), Inches(9.3), Inches(1.2), [
    ('Trust features carry no predictive power of their own — alone they hurt.', 15.5, False, INK2),
    ('Their entire value is modulating how much to trust the satellite, basin by basin.', 15.5, True, INK),
], spacing=1.35)
footer(s, 11)

# ================= 12 · THE HARD TEST (TRB) ===============================
s = slide()
header(s, 'the hard test', 'Hold out an entire region — and the real bottleneck appears')
s.shapes.add_picture(os.path.join(A, 'fig_trb.png'),
                     Inches(0.62), Inches(1.95), width=Inches(8.0))
text(s, Inches(8.95), Inches(2.05), Inches(3.9), Inches(4.3), [
    ('All 101 Tennessee basins', 16, True, INK),
    ('held out at once.', 16, True, INK),
    ('', 8, False, INK),
    ('The satellite prior stays good (0.68).', 13.5, False, INK2),
    ('', 6, False, INK),
    ('Blind calibrator: its corrections were', 13.5, False, INK2),
    ('fit to residuals it cannot explain —', 13.5, False, INK2),
    ('they mistransfer under shift (−0.13).', 13.5, False, INK2),
    ('', 6, False, INK),
    ('Give it the same embedding:', 13.5, False, INK2),
    ('calibration heals and the full model', 13.5, False, INK2),
    ('beats the baseline (0.707 vs 0.688).', 13.5, True, INK),
], spacing=1.22)
text(s, Inches(0.8), Inches(6.55), Inches(11.8), Inches(0.5),
     'The calibrator must be conditioned on everything the prior consumed.',
     size=16, bold=True, color=ACCENT)
footer(s, 12)

# ================= 13 · ROADMAP ===========================================
s = slide()
header(s, 'roadmap', 'From one fold to a paper')
ly = Inches(3.35)
ln = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1.1), ly, Inches(12.3), ly)
ln.line.color.rgb = LINE
ln.line.width = Pt(2.2)
steps = [
    ('22-fold cross-validation', 'full-conditioning configuration,\nall 618 basins → main results\ntable; 24 basins with no\nobservations excluded'),
    ('More regional hold-outs', 'repeat the Tennessee protocol\non other regions to stress-test\ncalibration transfer beyond\none case'),
    ('Scale to CONUS', '2,229 basins, all 18 regions;\nreconstruct genuinely missing\nTennessee records, validated\non daily aggregates'),
    ('Self-supervised encoder', 'only 618 of 9,012 basins have\nlabels — pretrain on all 9,012\nby masked reconstruction,\nthen fine-tune in Sat2Stream'),
]
xx = Inches(1.1)
step_w = Inches(2.85)
for i, (t, d) in enumerate(steps):
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL, xx + Inches(0.02), ly - Inches(0.09),
                             Inches(0.18), Inches(0.18))
    dot.fill.solid()
    dot.fill.fore_color.rgb = ACCENT
    dot.line.fill.background()
    text(s, xx - Inches(0.1), ly - Inches(1.05), step_w, Inches(0.8),
         f'{i+1} · {t}', size=16.5, bold=True, color=INK)
    text(s, xx - Inches(0.1), ly + Inches(0.25), step_w, Inches(1.6),
         d.split('\n'), size=13, color=INK2, spacing=1.2)
    xx += Inches(2.95)
text(s, Inches(1.1), Inches(6.0), Inches(11.2), Inches(0.9), [
    ('Target: KDD / AAAI.  Evaluation: mask parts of well-observed records (ground truth kept for scoring,', 14, False, INK2),
    ('573 basins) — plus reconstructing records that are genuinely missing, validated indirectly.', 14, False, INK2),
], spacing=1.35)
footer(s, 13)

# ================= 14 · TAKEAWAYS =========================================
s = slide()
header(s, 'takeaways', 'What we know now')
msgs = [
    ('The satellite improves the prior in every setting tested.',
     'CONUS fold: median NSE 0.530 → 0.576.  Tennessee hold-out: prior 0.666 → 0.684.  Trust features cover karst, where the modality’s physics fail.'),
    ('Calibration must see what the prior saw.',
     'hiding the satellite from stage 2 collapsed regional transfer (0.684 → 0.552); restoring it flipped the outcome — 0.707, above the no-satellite baseline.'),
    ('A negative result became the main finding.',
     'the Tennessee failure, diagnosed hypothesis by hypothesis, produced the condition-completeness principle and the final architecture.'),
]
yy = Inches(2.05)
for i, (h, d) in enumerate(msgs):
    text(s, Inches(1.0), yy, Inches(0.6), Inches(0.6), str(i + 1),
         size=30, bold=True, color=ACCENT)
    text(s, Inches(1.8), yy + Inches(0.02), Inches(10.6), Inches(0.55), h,
         size=19, bold=True, color=INK)
    text(s, Inches(1.8), yy + Inches(0.55), Inches(10.6), Inches(0.5), d,
         size=14.5, color=INK2)
    yy += Inches(1.35)
text(s, Inches(1.0), Inches(6.35), Inches(11.3), Inches(0.6),
     'discussion — 22-fold plan · more regional hold-outs · what else must the condition contain?',
     size=15, bold=False, color=FAINT, align=PP_ALIGN.CENTER)
footer(s, 14)

prs.save(OUT)
print('saved', OUT)
