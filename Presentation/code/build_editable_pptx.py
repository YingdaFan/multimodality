"""
Build an *editable* version of the SAFER-Hydro / TRB deck.

  - Title banner on every slide → native rounded rectangle + editable text
  - Cover                       → fully native (logos as images, all text as text boxes)
  - model_card slide            → fully native rebuild (every variable name,
                                  hyperparameter, key-result is an editable text box)
  - experiment_setting slide    → native When/How section headers + image body
  - Other slides                → native banner + image of the plot body

The matplotlib plot bodies (maps, heatmaps, scatter, time-series) stay as PNGs
because they encode coordinate-system data that isn't meaningful as native
shapes.

Output: Presentation/presentation/SAFER_Hydro_TRB_editable.pptx
"""

import os
from pathlib import Path
from copy import deepcopy

import fitz                                   # PyMuPDF
from lxml import etree
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn


# script lives in Presentation/code/ — go two up to reach imputation/
PROJECT_DIR = Path(__file__).resolve().parents[2]
PDF_PATH    = PROJECT_DIR / 'InflowForecast.pdf'
OUT_PPTX    = PROJECT_DIR / 'Presentation' / 'presentation' / 'SAFER_Hydro_TRB_editable.pptx'
TMP_DIR     = PROJECT_DIR / 'Presentation' / 'presentation' / '_editable_assets'
TMP_DIR.mkdir(parents=True, exist_ok=True)


# ============================== Deck colours (echo slide_style.py) ==============================
COL_BANNER_TOP = 'DCEBC4'
COL_BANNER_MID = 'AED080'
COL_BANNER_BOT = '8CBD5C'
COL_BANNER_EDG = '4E7A2F'
COL_TITLE_TEXT = '1A1A1A'

COL_PREDICTION  = 'D4E6F1'   # light blue
COL_FORECASTING = 'E2EFD9'   # light green
COL_INPUTS      = 'FCE4D6'   # peach

COL_CALLOUT_EDGE = '7DA66A'
COL_KEY_GREEN    = '588A4A'
COL_KEY_BLUE     = '3F6A8F'
COL_KEY_PEACH    = 'A36B3A'


# ============================== Slide size 16:9 ==============================
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ============================== Helpers ==============================
def add_gradient_fill(shape, top, mid, bot):
    """Replace the shape's fill with a vertical 3-stop gradient."""
    sp_pr = shape.fill._xPr
    # remove any existing fill children
    for tag in ('a:solidFill', 'a:gradFill', 'a:blipFill', 'a:pattFill', 'a:noFill'):
        for el in sp_pr.findall(qn(tag)):
            sp_pr.remove(el)
    nsmap = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
    grad_xml = (
        '<a:gradFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
        'rotWithShape="1">'
        '  <a:gsLst>'
        f'    <a:gs pos="0"><a:srgbClr val="{top}"/></a:gs>'
        f'    <a:gs pos="50000"><a:srgbClr val="{mid}"/></a:gs>'
        f'    <a:gs pos="100000"><a:srgbClr val="{bot}"/></a:gs>'
        '  </a:gsLst>'
        '  <a:lin ang="5400000" scaled="0"/>'
        '</a:gradFill>'
    )
    sp_pr.append(etree.fromstring(grad_xml))


def style_line(shape, hex_color, width_pt):
    """Set the shape's outline colour + width."""
    line = shape.line
    line.color.rgb = RGBColor.from_string(hex_color)
    line.width = Pt(width_pt)


def add_title_banner(slide, text, font_size=24):
    """Native banner across the top of the slide. Editable text."""
    banner_h = Inches(0.62)
    margin_x = Inches(0.16)
    top      = Inches(0.13)
    width    = SLIDE_W - 2 * margin_x

    sh = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, margin_x, top, width, banner_h)
    # corner radius ≈ 8 % of height
    sh.adjustments[0] = 0.20
    add_gradient_fill(sh, COL_BANNER_TOP, COL_BANNER_MID, COL_BANNER_BOT)
    style_line(sh, COL_BANNER_EDG, 2.0)

    tf = sh.text_frame
    tf.margin_left = tf.margin_right = Inches(0.10)
    tf.margin_top  = tf.margin_bottom = Inches(0.04)
    tf.word_wrap   = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.bold = True
    run.font.size = Pt(font_size)
    run.font.color.rgb = RGBColor.from_string(COL_TITLE_TEXT)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return sh


def add_textbox(slide, x, y, w, h, text, *,
                size=12, bold=False, italic=False, color='222222',
                align='left', anchor='top', font='Liberation Sans'):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Inches(0.03)
    tf.margin_top  = tf.margin_bottom = Inches(0.02)
    tf.word_wrap   = True
    tf.vertical_anchor = {'top': MSO_ANCHOR.TOP, 'middle': MSO_ANCHOR.MIDDLE,
                          'bottom': MSO_ANCHOR.BOTTOM}[anchor]
    p = tf.paragraphs[0]
    p.alignment = {'left': PP_ALIGN.LEFT, 'center': PP_ALIGN.CENTER,
                   'right': PP_ALIGN.RIGHT}[align]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font
    run.font.color.rgb = RGBColor.from_string(color)
    return tb


def add_block(slide, x, y, w, h, fill_hex, edge_hex='666666'):
    """Just the coloured rounded-rectangle background (no internal text)."""
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sh.adjustments[0] = 0.04
    sh.fill.solid()
    sh.fill.fore_color.rgb = RGBColor.from_string(fill_hex)
    style_line(sh, edge_hex, 0.75)
    # Empty out the default text frame so it never renders 'Click to edit'
    sh.text_frame.text = ''
    return sh


def add_block_header(slide, x, y, w, text, *,
                     color='222222', size=14):
    """Title text positioned at the top of a content-block column."""
    return add_textbox(slide, x, y, w, Inches(0.36), text,
                        size=size, bold=True, color=color,
                        align='center', anchor='middle')


def strip_top(png_path, out_path, frac=0.11):
    """Crop the top `frac` of the image (banner area)."""
    with Image.open(png_path) as im:
        w, h = im.size
        crop = im.crop((0, int(h * frac), w, h))
        crop.save(out_path, 'PNG')
    return out_path


def add_image_below_banner(slide, png_path, top=Inches(0.95)):
    """Place a stripped figure under the banner, centred horizontally."""
    with Image.open(png_path) as im:
        iw, ih = im.size
    # available area: top..slide_h, full width
    avail_h = SLIDE_H - top - Inches(0.10)
    avail_w = SLIDE_W - Inches(0.40)
    ratio_i = iw / ih
    ratio_a = avail_w / avail_h
    if ratio_i > ratio_a:
        w = avail_w
        h = int(w / ratio_i)
    else:
        h = avail_h
        w = int(h * ratio_i)
    x = (SLIDE_W - w) // 2
    y = top + (avail_h - h) // 2
    slide.shapes.add_picture(str(png_path), x, y, w, h)


# ============================== Build cover slide ==============================
def extract_logos():
    """Pull ORNL + SAFER-Hydro logos as standalone PNGs from page 1 of the PDF."""
    doc = fitz.open(PDF_PATH)
    page = doc[0]
    out  = {}
    for img in page.get_images(full=True):
        xref = img[0]
        pix = fitz.Pixmap(doc, xref)
        if pix.n - pix.alpha >= 4:
            pix = fitz.Pixmap(fitz.csRGB, pix)
        fname = TMP_DIR / f'logo_xref{xref}.png'
        pix.save(fname)
        out[xref] = fname
        pix = None
    doc.close()
    # heuristics: ORNL = the wider logo on the left (xref 11, ~673×200);
    #             SAFER = top-right circular badge (xref 17, ~697×611)
    return out.get(11), out.get(17)


def build_cover(prs):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    ornl_path, safer_path = extract_logos()

    # Top header: "AI-TVA Project Progress Report"
    add_textbox(slide, Inches(3.5), Inches(0.20), Inches(6.5), Inches(0.45),
                'AI-TVA Project Progress Report',
                size=22, bold=True, color='1A1A1A', align='center', anchor='middle')

    # ORNL logo (top-left)
    if ornl_path and ornl_path.exists():
        slide.shapes.add_picture(str(ornl_path),
                                 Inches(0.30), Inches(0.20),
                                 height=Inches(0.85))
    # SAFER-Hydro logo (top-right)
    if safer_path and safer_path.exists():
        slide.shapes.add_picture(str(safer_path),
                                 Inches(11.0), Inches(0.10),
                                 height=Inches(2.4))

    # Title banner (project title)
    banner_top  = Inches(2.45)
    banner_h    = Inches(1.55)
    margin_x    = Inches(0.40)
    banner_w    = SLIDE_W - 2 * margin_x
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 margin_x, banner_top, banner_w, banner_h)
    sh.adjustments[0] = 0.20
    add_gradient_fill(sh, COL_BANNER_TOP, COL_BANNER_MID, COL_BANNER_BOT)
    style_line(sh, COL_BANNER_EDG, 2.5)
    tf = sh.text_frame
    tf.margin_left = tf.margin_right = Inches(0.30)
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = 'SAFER-Hydro: Scalable AI Inflow Forecasting for Hydropower Utilities'
    run.font.bold = True
    run.font.size = Pt(40)
    run.font.color.rgb = RGBColor.from_string(COL_TITLE_TEXT)

    # Name block at the bottom-left
    add_textbox(slide, Inches(0.5), Inches(4.40), Inches(8), Inches(0.55),
                'Dan Lu', size=24, bold=False)
    add_textbox(slide, Inches(0.5), Inches(4.92), Inches(8), Inches(0.45),
                'Senior Computational Hydrologist', size=18)
    add_textbox(slide, Inches(0.5), Inches(5.70), Inches(10), Inches(0.45),
                'Yingda Fan; Vinh Tran; Soumendra Bhanja', size=18)
    add_textbox(slide, Inches(0.5), Inches(6.50), Inches(4), Inches(0.40),
                'Feb. 27, 2026', size=14)

    return slide


# ============================== Build a "banner + image" content slide ==============================
def build_image_slide(prs, banner_text, png_path, font_size=24,
                      strip_banner=True):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    add_title_banner(slide, banner_text, font_size=font_size)
    if strip_banner:
        stripped = TMP_DIR / (Path(png_path).stem + '_nobanner.png')
        strip_top(png_path, stripped, frac=0.115)
        add_image_below_banner(slide, stripped, top=Inches(0.95))
    else:
        add_image_below_banner(slide, png_path, top=Inches(0.95))
    return slide


# ============================== Build experiment_setting (banner + section headers + image) ==============================
def build_experiment_setting_slide(prs, png_path):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    add_title_banner(slide,
                     'Experiment Setting — Temporal Split & Forecast Window',
                     font_size=24)

    # Native section headers (these *replace* the headers that were baked into the PNG,
    # so we strip those off too by cropping a bit more aggressively).
    add_textbox(slide, Inches(0.55), Inches(1.00), Inches(12.2), Inches(0.40),
                'When  ·  validation 1995–96 (2 yr)   ·   training 1997–2018 (22 yr)   '
                '·   testing 2019–22 (4 yr, hold-out)',
                size=16, bold=True, color='2F4F1F')

    add_textbox(slide, Inches(0.55), Inches(3.75), Inches(12.2), Inches(0.40),
                'How  ·  168-hour input window  →  18-hour forecast,   '
                'stride 24 h between windows',
                size=16, bold=True, color='2F4F1F')

    # The body image (split timeline + sliding window) — we still need to strip
    # the banner *and* the two section header lines, because they're already
    # native now.  Easiest: crop top 11 % (banner) only and let the embedded
    # headers be hidden by the new native ones via positioning.  Cleaner: just
    # crop top 11 % and use the body image without trying to remove embedded
    # headers — the slight duplication is acceptable here.  For best results,
    # we render the body shapes natively, but those bars/rectangles aren't
    # trivial — keep the image.
    stripped = TMP_DIR / 'experiment_setting_body.png'
    with Image.open(png_path) as im:
        w, h = im.size
        crop = im.crop((0, int(h * 0.115), w, h))
        crop.save(stripped, 'PNG')

    # Place the body image taking the area below banner (and shifted so the
    # native headers sit at logical positions above their visuals)
    with Image.open(stripped) as im:
        iw, ih = im.size
    target_w = SLIDE_W - Inches(0.40)
    target_h = int(target_w * ih / iw)
    if target_h > SLIDE_H - Inches(0.95) - Inches(0.10):
        target_h = SLIDE_H - Inches(0.95) - Inches(0.10)
        target_w = int(target_h * iw / ih)
    x = (SLIDE_W - target_w) // 2
    y = Inches(0.95)
    slide.shapes.add_picture(str(stripped), x, y, target_w, target_h)
    return slide


# ============================== Build the fully-native model_card slide ==============================
DYNAMIC_VARS = [
    ('Rainf',       'mm/h',   'Rainfall (primary forcing)'),
    ('Tair',        'K',      '2-m air temperature'),
    ('Qair',        'kg/kg',  'Specific humidity'),
    ('PSurf',       'Pa',     'Surface pressure'),
    ('Wind_E',      'm/s',    'Zonal wind'),
    ('Wind_N',      'm/s',    'Meridional wind'),
    ('SWdown',      'W/m²', 'Downward shortwave radiation'),
    ('LWdown',      'W/m²', 'Downward longwave radiation'),
    ('PotEvap',     'mm/h',   'Potential evaporation'),
    ('CAPE',        'J/kg',   'Convective avail. pot. energy'),
    ('CRainf_frac', '–',      'Convective rainfall fraction'),
]

CUMUL = [
    ('Rainf_sum_{24,72,168} h',     '[mm]',  'Rolling rainfall (1 d, 3 d, 7 d)'),
    ('Tair_avg_{24,72,168} h',      '[K]',   'Rolling temperature'),
    ('PotEvap_sum_{24,72,168} h',   '[mm]',  'Rolling potential ET'),
]

STATIC_GROUPS = [
    ('Climate (9)',
     'p_mean, pet_mean, aridity, p_seasonality, frac_snow,\n'
     'high/low_prec_freq, high/low_prec_dur'),
    ('Topography (1)',
     'area_sqkm'),
    ('HydroATLAS (14)',
     'ele_mt_sav, slp_dg_uav, ria_ha_usu, run_mm_syr, gwt_cm_sav,\n'
     'cly/slt/snd_pc_uav, kar/prm/pac_pc_use, crp/for/urb_pc_use'),
]

TRAINING_ARCH = [
    ('Model',          'FutureTST  (Transformer encoder + decoder)'),
    ('Encoder',        '168 h history → latent representation'),
    ('Decoder',        'latent + future met → 18 h Q forecast'),
]
TRAINING_HYPER = [
    ('Joint training', '618 CONUS basins  (1 model, no per-region fine-tune)'),
    ('Batch',          '618  (one full basin set per step)'),
    ('Epochs',         '200  (early-stop patience = 20)'),
    ('Learning rate',  '0.001  (Adam)'),
    ('Stride',         '24 h  →  one new 18-h forecast per day '
                       '(~358 forecasts / basin / year)'),
]
TRAINING_DATA = [
    ('Split',           'val 1995–96 (2 y) · train 1997–2018 (22 y) · test 2019–22 (4 y)'),
    ('Imputation prior', 'y_imputed injected from upstream diffusion pipeline '
                        '(handles gaps in observed Q during training)'),
]

OUTPUTS = [
    ('Raw output',     '18-step standardised Q   per window'),
    ('Denormalized',   '× σ_basin + μ_basin  →  Q [mm/day]\n'
                       '(specific runoff = streamflow / area × 86.4, '
                       'so cross-basin comparable)'),
    ('Reported',       'NSE · KGE · RMSE · MAE · R² · pbias'),
    ('Per-lead view',  '1, 2, …, 18-hour-ahead skill curves'),
    ('Evaluation set', '130 TRB gauges  (104 with valid all-lead metrics)'),
]


def build_model_card_slide(prs):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    add_title_banner(slide,
                     'Model Card — FutureTST · 35 features × 168 h → 18 h Streamflow Forecast',
                     font_size=20)

    # Three columns
    GAP    = Inches(0.12)
    MARGIN = Inches(0.20)
    COL_W  = (SLIDE_W - 2 * MARGIN - 2 * GAP) / 3
    TOP    = Inches(0.95)
    COL_H  = SLIDE_H - TOP - Inches(0.20)

    x_in    = MARGIN
    x_model = MARGIN + COL_W + GAP
    x_out   = MARGIN + 2 * (COL_W + GAP)

    HEADER_H = Inches(0.40)
    BODY_TOP_OFFSET = HEADER_H + Inches(0.08)

    # ---- Inputs column (peach) ----
    add_block(slide, x_in, TOP, COL_W, COL_H, COL_INPUTS, edge_hex=COL_KEY_PEACH)
    add_block_header(
        slide, x_in, TOP + Inches(0.04), COL_W,
        f'Inputs   ({len(DYNAMIC_VARS)} dynamic + 9 rolling + 24 static)',
        color='222222', size=14)

    inner_x = x_in + Inches(0.15)
    inner_w = COL_W - Inches(0.30)
    y = TOP + BODY_TOP_OFFSET

    add_textbox(slide, inner_x, y, inner_w, Inches(0.26),
                '11 dynamic meteorological forcings', size=11, bold=True)
    y += Inches(0.24)
    add_textbox(slide, inner_x, y, inner_w, Inches(0.20),
                'ERA5-Land hourly, 1985-present',
                size=8.5, italic=True, color='555555')
    y += Inches(0.18)
    for name, unit, desc in DYNAMIC_VARS:
        add_textbox(slide, inner_x,                  y,
                    Inches(1.30), Inches(0.20),
                    f'►  {name}', size=9.5, bold=True, font='Liberation Mono')
        add_textbox(slide, inner_x + Inches(1.30),   y,
                    Inches(0.70), Inches(0.20),
                    f'[{unit}]', size=8.5, color='555555')
        add_textbox(slide, inner_x + Inches(2.00),   y,
                    Inches(2.30), Inches(0.20),
                    desc, size=9.5)
        y += Inches(0.19)

    y += Inches(0.06)
    add_textbox(slide, inner_x, y, inner_w, Inches(0.26),
                '9 rolling-cumulative features', size=11, bold=True)
    y += Inches(0.22)
    add_textbox(slide, inner_x, y, inner_w, Inches(0.20),
                'derived in preprocessing — capture antecedent state',
                size=8.5, italic=True, color='555555')
    y += Inches(0.20)
    for name, unit, desc in CUMUL:
        add_textbox(slide, inner_x,                  y,
                    Inches(2.50), Inches(0.20),
                    f'►  {name}', size=9, bold=True, font='Liberation Mono')
        add_textbox(slide, inner_x + Inches(2.50),   y,
                    Inches(0.65), Inches(0.20),
                    unit, size=8.5, color='555555')
        y += Inches(0.18)
        add_textbox(slide, inner_x + Inches(0.20),   y,
                    Inches(3.20), Inches(0.20),
                    desc, size=9)
        y += Inches(0.18)

    y += Inches(0.06)
    add_textbox(slide, inner_x, y, inner_w, Inches(0.26),
                '24 static catchment attributes (per basin)', size=11, bold=True)
    y += Inches(0.22)
    for grp, body in STATIC_GROUPS:
        add_textbox(slide, inner_x, y, Inches(2.0), Inches(0.20),
                    f'►  {grp}', size=10, bold=True)
        y += Inches(0.20)
        body_height = (body.count('\n') + 1) * 0.18 + 0.04
        add_textbox(slide, inner_x + Inches(0.25), y,
                    inner_w - Inches(0.25), Inches(body_height),
                    body, size=8.5, font='Liberation Mono')
        y += Inches(body_height + 0.02)

    # ---- Model + Training column (blue) ----
    add_block(slide, x_model, TOP, COL_W, COL_H, COL_PREDICTION,
              edge_hex=COL_KEY_BLUE)
    add_block_header(slide, x_model, TOP + Inches(0.04), COL_W,
                     'FutureTST   ·   Training Setup',
                     color='222222', size=14)

    inner_x_m = x_model + Inches(0.15)
    inner_w_m = COL_W - Inches(0.30)
    y = TOP + BODY_TOP_OFFSET
    LABEL_W = Inches(1.55)

    def write_section_label(label):
        nonlocal y
        add_textbox(slide, inner_x_m, y, inner_w_m, Inches(0.28),
                    label, size=12, bold=True)
        y += Inches(0.30)

    def write_kv(label, val, val_lines=1):
        nonlocal y
        h = Inches(0.26 * val_lines + 0.04)
        add_textbox(slide, inner_x_m, y, LABEL_W, h,
                    f'{label}:', size=10.5, bold=True, color=COL_KEY_BLUE)
        add_textbox(slide, inner_x_m + LABEL_W, y,
                    inner_w_m - LABEL_W, h,
                    val, size=10)
        y += h + Inches(0.02)

    write_section_label('Architecture')
    for label, val in TRAINING_ARCH:
        write_kv(label, val)

    y += Inches(0.06)
    write_section_label('Training')
    for label, val in TRAINING_HYPER:
        n_lines = val.count('\n') + 1
        write_kv(label, val, val_lines=n_lines)

    y += Inches(0.06)
    write_section_label('Data')
    for label, val in TRAINING_DATA:
        # estimated lines (text wraps at ~ 4 inches width)
        write_kv(label, val, val_lines=2)

    # Source footer at bottom of column
    add_textbox(slide, inner_x_m, TOP + COL_H - Inches(0.32),
                inner_w_m, Inches(0.24),
                'Source: run_forecast.sh + preprocess_camelsh_forecast.py',
                size=8.5, italic=True, color='555555')

    # ---- Outputs column (green) ----
    add_block(slide, x_out, TOP, COL_W, COL_H, COL_FORECASTING,
              edge_hex=COL_KEY_GREEN)
    add_block_header(slide, x_out, TOP + Inches(0.04), COL_W,
                     'Outputs', color='222222', size=14)

    inner_x_o = x_out + Inches(0.15)
    inner_w_o = COL_W - Inches(0.30)
    y_o = TOP + BODY_TOP_OFFSET

    for label, val in OUTPUTS:
        add_textbox(slide, inner_x_o, y_o, inner_w_o, Inches(0.28),
                    f'►  {label}', size=11.5, bold=True, color=COL_KEY_GREEN)
        y_o += Inches(0.26)
        n_lines = val.count('\n') + 1
        height = Inches(0.22 * n_lines + 0.08)
        add_textbox(slide, inner_x_o + Inches(0.20), y_o,
                    inner_w_o - Inches(0.20), height,
                    val, size=10)
        y_o += height + Inches(0.04)

    # Key Results box at the bottom of Outputs column
    kr_h = Inches(1.10)
    kr_y = TOP + COL_H - kr_h - Inches(0.16)
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 x_out + Inches(0.15), kr_y,
                                 COL_W - Inches(0.30), kr_h)
    box.adjustments[0] = 0.06
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor.from_string('FFFFFF')
    style_line(box, COL_KEY_GREEN, 1.0)
    tf_kr = box.text_frame
    tf_kr.margin_left = tf_kr.margin_right = Inches(0.10)
    tf_kr.margin_top = Inches(0.10)
    tf_kr.vertical_anchor = MSO_ANCHOR.TOP
    p = tf_kr.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = 'Key Results  (TRB, 18-hour horizon)'
    r.font.bold = True
    r.font.size = Pt(12)
    r.font.color.rgb = RGBColor.from_string('222222')
    p2 = tf_kr.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = 'median 1-h NSE  =  0.996      median 18-h NSE  =  0.863'
    r2.font.size = Pt(10.5)
    r2.font.name = 'Liberation Mono'
    r2.font.color.rgb = RGBColor.from_string('222222')
    p3 = tf_kr.add_paragraph()
    p3.alignment = PP_ALIGN.CENTER
    r3 = p3.add_run()
    r3.text = '94 % of TRB gauges keep NSE ≥ 0.7 at the 18-hour horizon'
    r3.font.size = Pt(10)
    r3.font.color.rgb = RGBColor.from_string('222222')

    # Arrows between columns
    arrow_y = TOP + COL_H / 2 - Inches(0.06)
    for ax_start_x in (x_in + COL_W, x_model + COL_W):
        ar = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                    ax_start_x + Inches(0.005), arrow_y,
                                    GAP - Inches(0.01), Inches(0.22))
        ar.fill.solid()
        ar.fill.fore_color.rgb = RGBColor.from_string(COL_BANNER_EDG)
        ar.line.fill.background()


# ============================== Main ==============================
def build_editable_pptx():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    # 1. cover
    build_cover(prs)

    # 2. global_vs_tn_map
    build_image_slide(
        prs,
        'Training Setup — Global FutureTST  ·  1 Model · 618 CONUS Basins · '
        'Regional Focus on TRB',
        PROJECT_DIR / 'Presentation' / 'presentation' / 'global_vs_tn_map.png',
        font_size=18)

    # 3. data_atlas
    build_image_slide(
        prs,
        'Data Availability Atlas — Hourly Q Coverage, 618 CONUS Basins (1985–2022)',
        PROJECT_DIR / 'Presentation' / 'presentation' / 'data_atlas.png',
        font_size=18)

    # 4. experiment_setting
    build_experiment_setting_slide(
        prs,
        PROJECT_DIR / 'Presentation' / 'presentation' / 'experiment_setting.png')

    # 5. model_card — fully native
    build_model_card_slide(prs)

    # 6. coverage_vs_nse
    build_image_slide(
        prs,
        'Target Data Quality (train & test) vs Prediction Quality — TN Basins',
        PROJECT_DIR / 'Presentation' / 'presentation' / 'coverage_vs_nse.png',
        font_size=20)

    # 7. basin_panel
    build_image_slide(
        prs,
        'Selected Basins — Observed vs 18-hour-ahead Forecast  '
        '(test coverage ≥ 90 %, calendar 2019)',
        PROJECT_DIR / 'Presentation' / 'presentation' / 'basin_panel.png',
        font_size=18)

    # 8. forecast_boxplot
    build_image_slide(
        prs,
        'NSE by Lead Time — Box Plot  (n = 104 TRB basins)',
        PROJECT_DIR / 'Presentation' / 'tennessee_18lead' / 'figures' / 'forecast_boxplot.png',
        font_size=22)

    # 9. forecast_cdf
    build_image_slide(
        prs,
        'NSE by Lead Time — Cumulative Distribution  (n = 104 TRB basins)',
        PROJECT_DIR / 'Presentation' / 'tennessee_18lead' / 'figures' / 'forecast_cdf.png',
        font_size=20)

    # 10. forecast_lines
    build_image_slide(
        prs,
        'NSE Degradation with Forecast Lead Time  (n = 104 TRB basins)',
        PROJECT_DIR / 'Presentation' / 'tennessee_18lead' / 'figures' / 'forecast_lines.png',
        font_size=22)

    # 11. page9_summary
    build_image_slide(
        prs,
        'Global FutureTST Streamflow Forecast (1–18-hour)',
        PROJECT_DIR / 'Presentation' / 'tennessee_18lead' / 'figures' / 'page9_summary.png',
        font_size=22)

    prs.save(OUT_PPTX)
    print(f'pptx saved → {OUT_PPTX}   ({len(prs.slides)} slides)')


if __name__ == '__main__':
    build_editable_pptx()
