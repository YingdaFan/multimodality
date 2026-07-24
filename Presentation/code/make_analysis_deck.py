"""
Build  Presentation/analyze.pptx  —  a focused diagnostic deck answering:

    "Which TN basins forecast poorly at long lead (15/16/17/18 h),
     and WHY?  training-data scarcity / test-data scarcity / hydrology?"

Inputs (all already on disk):
  - Presentation/tennessee_18lead/tn_nse_wide.csv      per-basin NSE, every 1..18 h lead
  - Presentation/tennessee_18lead/tn_metrics_long.csv  n_valid (test obs count) per lead
  - Presentation/presentation/tn_target_coverage.csv   train_cov / test_cov per basin
  - ../camelsh_tennessee.parquet                       static catchment attributes
  - Presentation/tennessee_18lead/figures/timeseries_*_jan2019.png   (flashiness illustration)

Outputs:
  - Presentation/tennessee_18lead/analysis/*.png       generated charts
  - Presentation/tennessee_18lead/late_lead_analysis.csv  master table
  - Presentation/analyze.pptx
"""

import os as _os
from pathlib import Path as _Path
_os.chdir(_Path(__file__).resolve().parents[2])      # -> imputation/

import os
import numpy as np
import pandas as pd
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
from matplotlib.lines import Line2D

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ----------------------------------------------------------------------------
HERE      = 'Presentation'
WIDE_CSV  = f'{HERE}/tennessee_18lead/tn_nse_wide.csv'
LONG_CSV  = f'{HERE}/tennessee_18lead/tn_metrics_long.csv'
COV_CSV   = f'{HERE}/presentation/tn_target_coverage.csv'
PARQUET   = '../camelsh_tennessee.parquet'
TS_DIR    = f'{HERE}/tennessee_18lead/figures'
FIG_DIR   = f'{HERE}/tennessee_18lead/analysis'
PPTX_OUT  = f'{HERE}/analyze.pptx'
os.makedirs(FIG_DIR, exist_ok=True)

EXCLUDED  = {'03566420'}
THR       = 0.70                      # "poor" = mean NSE over 15-18h below this
LATE      = ['15h_ahead', '16h_ahead', '17h_ahead', '18h_ahead']

# deck palette (echoes slide_style.py)
GREEN_DK  = RGBColor(0x4e, 0x7a, 0x2f)
GREEN_BAN = RGBColor(0x8c, 0xbd, 0x5c)
SAGE      = RGBColor(0xae, 0xca, 0xa1)
INK       = RGBColor(0x1a, 0x1a, 0x1a)
GREY      = RGBColor(0x55, 0x55, 0x55)
RED       = RGBColor(0xc0, 0x39, 0x2b)
WHITE     = RGBColor(0xff, 0xff, 0xff)

C_TRAIN   = '#d1495b'     # training-data scarce
C_BOTH    = '#8e44ad'     # train + test scarce
C_HYDRO   = '#2e86c1'     # hydrology / features
C_GOOD    = '#c9d6c2'     # the healthy 95 basins

plt.rcParams.update({
    'font.family': 'Liberation Sans', 'font.size': 13,
    'axes.titlesize': 16, 'axes.titleweight': 'bold',
    'axes.labelsize': 14, 'axes.edgecolor': '#444', 'axes.linewidth': 0.9,
    'figure.dpi': 150,
})

# ============================ assemble master table ==========================
w = pd.read_csv(WIDE_CSV, dtype={'basin': str})
allcols = [f'{h}h_ahead' for h in range(1, 19)]
w = w.dropna(subset=allcols)
w = w[~w['basin'].isin(EXCLUDED)].copy()
w['late_mean'] = w[LATE].mean(axis=1)
m = w.set_index('basin')

cov = pd.read_csv(COV_CSV, dtype={'basin_id': str}).set_index('basin_id')
m = m.join(cov[['train_cov', 'test_cov']])

ml = pd.read_csv(LONG_CSV, dtype={'basin': str})
m['n_valid_test'] = ml[ml['hours_ahead'] == 18].set_index('basin')['n_valid']

sa = pd.read_parquet(PARQUET,
        columns=['basin_id', 'area_sqkm', 'baseflow_index', 'slope_pct',
                 'elev_mean', 'high_prec_dur']).groupby('basin_id').first()
m = m.join(sa)

m.to_csv(f'{HERE}/tennessee_18lead/late_lead_analysis.csv')

# 03491544 is a SHORT-lead anomaly, not a long-lead failure: 1h NSE 0.47 < 18h 0.83.
# Its 15-18h mean dips below THR only because the mid-leads are weak, so it would
# mislead a long-lead-failure story. Set it aside (reported as a noted exception).
ANOM = '03491544'
ANALY = m.drop(index=[ANOM]) if ANOM in m.index else m       # 104 -> 103 analysed
poor = ANALY[ANALY['late_mean'] < THR].sort_values('late_mean')
good = ANALY[ANALY['late_mean'] >= THR]
POOR = list(poor.index)
n_all, n_poor = len(ANALY), len(poor)
anom_row = m.loc[ANOM] if ANOM in m.index else None
print(f'{n_poor} poor / {n_all} basins  (threshold mean NSE_15-18h < {THR}; '
      f'1 anomaly {ANOM} set aside)')

# --- cause label per poor basin --------------------------------------------
def cause(r):
    train_bad = r['train_cov'] < 0.50
    test_bad  = (r['test_cov'] < 0.50) or (r['n_valid_test'] < 300)
    if train_bad and test_bad: return 'both'
    if train_bad:              return 'train'
    return 'hydro'

poor = poor.copy()
poor['cause'] = poor.apply(cause, axis=1)
CAUSE_C = {'train': C_TRAIN, 'both': C_BOTH, 'hydro': C_HYDRO}
CAUSE_LBL = {'train': 'Training data scarce',
             'both':  'Train + test data scarce',
             'hydro': 'Hydrology (flashy) — data OK'}

# core: in the bottom-5 at every one of 15/16/17/18h (computed on analysed set)
core = set.intersection(*[set(ANALY[c].nsmallest(5).index) for c in LATE])

# ============================ figures =======================================
def _despine(ax):
    for s in ('top', 'right'):
        ax.spines[s].set_visible(False)

# fig A — degradation curves 1..18h for the poor basins
def fig_degrade():
    fig, ax = plt.subplots(figsize=(11, 6.2))
    xs = list(range(1, 19))
    # healthy reference band
    band = m.loc[good.index, allcols]
    ax.fill_between(xs, band.quantile(.25), band.quantile(.75),
                    color=C_GOOD, alpha=.7, label='Healthy basins (IQR)', zorder=1)
    ax.plot(xs, band.median(), color='#5f7a52', lw=2, ls='--',
            label='Healthy median', zorder=2)
    for b in POOR:
        c = CAUSE_C[poor.loc[b, 'cause']]
        ax.plot(xs, m.loc[b, allcols].values, color=c, lw=2, alpha=.9,
                marker='o', ms=3, zorder=3)
    # de-cluttered end-labels: stagger vertically so the 0.6-0.7 pile-up is legible
    ends = sorted((m.loc[b, '18h_ahead'], b) for b in POOR)
    last = -9.0
    for yv, b in ends:
        ylab = max(yv, last + 0.030)
        last = ylab
        c = CAUSE_C[poor.loc[b, 'cause']]
        ax.annotate(b, xy=(18, yv), xytext=(18.8, ylab), textcoords='data',
                    fontsize=8.5, color=c, va='center', fontweight='bold',
                    arrowprops=dict(arrowstyle='-', color=c, lw=.5, alpha=.5))
    ax.axvspan(14.5, 18.5, color='#f2c14e', alpha=.18, zorder=0)
    ax.text(16.5, ax.get_ylim()[0] + .03, '15–18 h\nfocus', ha='center',
            fontsize=10, color='#8a6d1a', fontweight='bold')
    ax.set_xlabel('Forecast lead time (hours)')
    ax.set_ylabel('NSE')
    ax.set_xticks(xs)
    ax.set_xlim(.5, 19.8)
    ax.grid(axis='y', alpha=.25)
    _despine(ax)
    handles = [Line2D([], [], color=C_TRAIN, lw=2, label='Training scarce'),
               Line2D([], [], color=C_BOTH,  lw=2, label='Train+test scarce'),
               Line2D([], [], color=C_HYDRO, lw=2, label='Hydrology (data OK)')]
    ax.legend(handles=handles + [
        Line2D([], [], color='#5f7a52', lw=2, ls='--', label='Healthy median')],
        loc='lower left', framealpha=.9, fontsize=11)
    fig.tight_layout()
    p = f'{FIG_DIR}/degrade.png'; fig.savefig(p, dpi=170); plt.close(); return p

# generic scatter: late NSE vs a driver, poor basins highlighted
def fig_scatter(xcol, xlabel, fname, logx=False, vline=None, color_by_cause=True):
    fig, ax = plt.subplots(figsize=(9.6, 6.4))
    ax.scatter(good[xcol], good['late_mean'], s=55, c=C_GOOD,
               edgecolor='#8a9a82', linewidth=.5, label='Healthy basins', zorder=2)
    for b in POOR:
        c = CAUSE_C[poor.loc[b, 'cause']] if color_by_cause else RED.__str__()
        ax.scatter(m.loc[b, xcol], m.loc[b, 'late_mean'], s=130, c=c,
                   edgecolor='black', linewidth=.8, zorder=4)
        ax.annotate(b, (m.loc[b, xcol], m.loc[b, 'late_mean']),
                    fontsize=9, xytext=(5, 4), textcoords='offset points')
    if logx:
        ax.set_xscale('log')
    if vline is not None:
        ax.axvline(vline, color='gray', ls=':', lw=1.3, alpha=.8)
    r = float(ANALY[['late_mean', xcol]].corr().iloc[0, 1])
    ax.set_xlabel(xlabel)
    ax.set_ylabel('Mean NSE  (15–18 h leads)')
    ax.axhline(THR, color='#c0392b', ls=':', lw=1.2, alpha=.7)
    ax.text(0.015, 0.04, f'Pearson r = {r:+.2f}   (n = {n_all})',
            transform=ax.transAxes, fontsize=12, fontweight='bold',
            bbox=dict(boxstyle='round,pad=.3', fc='white', ec='#888', lw=.6))
    ax.grid(True, alpha=.25)
    _despine(ax)
    fig.tight_layout()
    p = f'{FIG_DIR}/{fname}'; fig.savefig(p, dpi=170); plt.close(); return p

# fig features: area (log) vs NSE colored by BFI
def fig_features():
    fig, ax = plt.subplots(figsize=(9.6, 6.4))
    sc = ax.scatter(good['area_sqkm'], good['late_mean'], c=good['baseflow_index'],
                    cmap='YlGnBu', s=70, edgecolor='#555', linewidth=.4,
                    vmin=20, vmax=75, zorder=2)
    cb = plt.colorbar(sc, ax=ax, pad=.02); cb.set_label('Baseflow index')
    for b in POOR:
        ax.annotate(b, (m.loc[b, 'area_sqkm'], m.loc[b, 'late_mean']),
                    fontsize=9, xytext=(5, 4), textcoords='offset points',
                    color=CAUSE_C[poor.loc[b, 'cause']], fontweight='bold')
        ax.scatter(m.loc[b, 'area_sqkm'], m.loc[b, 'late_mean'], s=150,
                   facecolor='none', edgecolor=CAUSE_C[poor.loc[b, 'cause']],
                   linewidth=1.8, zorder=4)
    ax.set_xscale('log')
    ax.set_xlabel('Catchment area  (km², log scale)')
    ax.set_ylabel('Mean NSE  (15–18 h leads)')
    ax.axhline(THR, color='#c0392b', ls=':', lw=1.2, alpha=.7)
    ax.grid(True, alpha=.25)
    _despine(ax)
    fig.tight_layout()
    p = f'{FIG_DIR}/features.png'; fig.savefig(p, dpi=170); plt.close(); return p

# repeat offenders = in the worst-6 at EVERY one of 15/16/17/18 h
REPEAT = [b for b in POOR if b in core]              # POOR already worst-first

def rgb255(c):
    return tuple(int(x * 255) for x in matplotlib.colors.to_rgb(c))

# CONSISTENT COLOUR LANGUAGE (whole deck): colour = root cause; box = repeat offender
def cause_color(b):
    return CAUSE_C[poor.loc[b, 'cause']] if b in poor.index else '#7a7a7a'

# fig per-lead: the worst-6 basins at each of 15/16/17/18 h, with the
# recurring (repeat-offender) basins boxed + linked across columns,
# coloured by ROOT CAUSE so the colour means the same thing everywhere.
def fig_perlead(nshow=6):
    fig, ax = plt.subplots(figsize=(10.8, 6.3))
    pos = {}
    for ci, c in enumerate(LATE):
        for rank, (b, nse) in enumerate(ANALY[c].nsmallest(nshow).items()):
            y = nshow - rank                          # worst at top
            pos[(b, ci)] = y
            rep = b in core
            col = cause_color(b) if rep else '#7a7a7a'
            ax.text(ci, y, f'{b}\n{nse:.2f}', ha='center', va='center',
                    fontsize=10.5, color='white' if rep else col,
                    fontweight='bold' if rep else 'normal', zorder=3,
                    bbox=dict(boxstyle='round,pad=0.3',
                              fc=col if rep else 'white',
                              ec=col, lw=1.3, alpha=0.95 if rep else 0.0))
    for b in core:                                    # link recurring basins
        xs = [ci for ci in range(len(LATE)) if (b, ci) in pos]
        ax.plot(xs, [pos[(b, ci)] for ci in xs],
                color=cause_color(b), lw=2, alpha=0.40, zorder=1)
    ax.set_xticks(range(len(LATE)))
    ax.set_xticklabels(['15 h', '16 h', '17 h', '18 h'], fontsize=14, fontweight='bold')
    ax.set_xlim(-0.5, len(LATE) - 0.5)
    ax.set_ylim(0.3, nshow + 1.4)
    ax.set_yticks([])
    ax.set_xlabel('Forecast lead time', fontsize=14)
    ax.text(-0.48, nshow + 1.05, 'worst (top of column)', fontsize=10.5,
            color='#999', style='italic')
    for sp in ('top', 'right', 'left'):
        ax.spines[sp].set_visible(False)
    # cause legend (colour = root cause, explained on slides 3-5)
    handles = [Line2D([], [], marker='s', ls='', ms=11, mfc=C_TRAIN, mec='none',
                      label='Training scarce'),
               Line2D([], [], marker='s', ls='', ms=11, mfc=C_BOTH, mec='none',
                      label='Train+test scarce'),
               Line2D([], [], marker='s', ls='', ms=11, mfc=C_HYDRO, mec='none',
                      label='Hydrology (data OK)'),
               Line2D([], [], marker='s', ls='', ms=11, mfc='white', mec='#7a7a7a',
                      label='Poor at this lead only')]
    ax.legend(handles=handles, loc='upper center', ncol=4, fontsize=9.5,
              frameon=False, bbox_to_anchor=(0.5, 1.005),
              handletextpad=0.3, columnspacing=1.0,
              title='box colour = root cause (slides 3–5)', title_fontsize=9.5)
    fig.tight_layout()
    p = f'{FIG_DIR}/perlead.png'; fig.savefig(p, dpi=170); plt.close(); return p

print('rendering figures...')
P_PERLEAD  = fig_perlead()
P_DEGRADE  = fig_degrade()
P_TRAIN    = fig_scatter('train_cov', 'Training-period target coverage  (1997–2018)',
                         'nse_vs_train.png', vline=0.50)
P_TEST     = fig_scatter('test_cov',  'Test-period target coverage  (2019–2022)',
                         'nse_vs_test.png', vline=0.50)
P_FEATURES = fig_features()

# ============================ PPTX ==========================================
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height

def add_slide():
    return prs.slides.add_slide(BLANK)

def banner(slide, title, page):
    bar = slide.shapes.add_shape(1, 0, 0, SW, Inches(0.95))
    bar.fill.solid(); bar.fill.fore_color.rgb = GREEN_BAN
    bar.line.color.rgb = GREEN_DK; bar.line.width = Pt(1)
    tf = bar.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.35); tf.margin_top = Inches(0.1)
    p = tf.paragraphs[0]; p.text = title
    p.font.size = Pt(26); p.font.bold = True; p.font.color.rgb = INK
    p.font.name = 'Liberation Sans'
    badge = slide.shapes.add_shape(9, SW - Inches(0.95), Inches(0.16),
                                   Inches(0.63), Inches(0.63))
    badge.fill.solid(); badge.fill.fore_color.rgb = WHITE
    badge.line.color.rgb = GREEN_DK
    bp = badge.text_frame.paragraphs[0]; bp.text = str(page)
    bp.font.size = Pt(18); bp.font.bold = True; bp.font.color.rgb = GREEN_DK
    bp.alignment = PP_ALIGN.CENTER

def textbox(slide, left, top, width, height, lines, size=14, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True
    for i, (txt, opts) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = txt; p.alignment = align
        p.font.size = Pt(opts.get('size', size))
        p.font.bold = opts.get('bold', False)
        p.font.color.rgb = opts.get('color', INK)
        p.font.name = 'Liberation Sans'
        if opts.get('space_before'):
            p.space_before = Pt(opts['space_before'])
    return tb

def callout(slide, text, top=None):
    top = top if top is not None else SH - Inches(1.05)
    box = slide.shapes.add_shape(5, Inches(0.4), top, SW - Inches(0.8), Inches(0.8))
    box.fill.solid(); box.fill.fore_color.rgb = SAGE
    box.line.color.rgb = GREEN_DK; box.line.width = Pt(1)
    tf = box.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.3)
    p = tf.paragraphs[0]; p.text = '❖  ' + text
    p.font.size = Pt(15); p.font.bold = True; p.font.color.rgb = INK
    p.font.name = 'Liberation Sans'

def pic(slide, path, left, top, height=None, width=None):
    if height is not None:
        return slide.shapes.add_picture(path, left, top, height=height)
    return slide.shapes.add_picture(path, left, top, width=width)

# ---- Slide 1 : worst basins at each lead --------------------------------
s = add_slide()
banner(s, 'Worst basins at each lead — 15 / 16 / 17 / 18 h', 1)
pic(s, P_PERLEAD, Inches(0.25), Inches(1.1), height=Inches(5.1))
lines = [('Each column = the 6 lowest-NSE basins', {'bold': True, 'size': 13}),
         ('at that lead time.', {'bold': True, 'size': 13}),
         ('Coloured + boxed basins appear in', {'size': 12.5, 'space_before': 10}),
         ('every column — they recur at all four', {'size': 12.5}),
         ('lead times.', {'size': 12.5}),
         (f'-> {len(REPEAT)} "repeat offenders":', {'bold': True, 'size': 13, 'color': RED, 'space_before': 12})]
for b in REPEAT:
    lines.append((f'{b}', {'bold': True, 'size': 13.5, 'color': RGBColor(*rgb255(cause_color(b))),
                           'space_before': 2}))
lines.append(('Box colour already encodes the root', {'size': 11, 'color': GREY, 'space_before': 12}))
lines.append(('cause — explained on slides 3–5.', {'size': 11, 'color': GREY}))
textbox(s, Inches(9.6), Inches(1.3), Inches(3.5), Inches(5.2), lines)
callout(s, f'The same {len(REPEAT)} basins sit in the worst-6 at every lead from 15 to 18 h '
           f'— these are the "repeat offenders" we dig into next.')

# ---- Slide 2 : repeat offenders + full poor set -------------------------
s = add_slide()
banner(s, 'Repeat offenders & the full poor set', 2)
pic(s, P_DEGRADE, Inches(0.25), Inches(1.15), height=Inches(5.0))
lines = [('Repeat offenders  (worst at all 4 leads)', {'bold': True, 'size': 13.5, 'color': RED})]
for b in REPEAT:
    lines.append((f'{b}   mean 15–18h NSE = {m.loc[b,"late_mean"]:.2f}',
                  {'size': 12.5, 'color': RGBColor(*rgb255(cause_color(b)))}))
lines.append((f'Other poor basins  (mean < {THR:.2f})',
              {'bold': True, 'size': 13.5, 'space_before': 12}))
for b in [x for x in POOR if x not in core]:
    lines.append((f'{b}   {m.loc[b,"late_mean"]:.2f}', {'size': 12.5, 'color': GREY}))
lines.append(('All 4 repeat offenders also rank among', {'size': 11.5, 'color': INK, 'space_before': 12}))
lines.append((f'these {n_poor} lowest-mean basins (cut-off NSE {THR:.2f}).', {'size': 11.5, 'color': INK}))
if anom_row is not None:
    lines.append((f'Set aside: {ANOM} — short-lead anomaly', {'size': 10.5, 'color': GREY, 'space_before': 8}))
    lines.append((f'(1h NSE {anom_row["1h_ahead"]:.2f} < 18h {anom_row["18h_ahead"]:.2f}),',
                  {'size': 10.5, 'color': GREY}))
    lines.append(('not a long-lead failure.', {'size': 10.5, 'color': GREY}))
textbox(s, Inches(9.5), Inches(1.2), Inches(3.6), Inches(5.4), lines)
callout(s, f'{n_poor} of {n_all} basins fall below NSE {THR:.2f} at long lead; '
           f'the {len(REPEAT)} repeat offenders are the worst. Why? -> data, or hydrology?')

# ---- Slide 3 : training-data scarcity ------------------------------------
s = add_slide()
banner(s, 'Cause 1 — Training-data scarcity (strongest)', 3)
pic(s, P_TRAIN, Inches(0.3), Inches(1.15), height=Inches(5.05))
med_pt = poor['train_cov'].median(); med_gd = good['train_cov'].median()
nlow = int((poor['train_cov'] < 0.5).sum())
lines = [('Finding', {'bold': True, 'size': 16, 'color': GREEN_DK}),
         (f'Poor basins median train coverage = {med_pt:.2f}', {'size': 14, 'space_before': 6}),
         (f'Healthy basins median = {med_gd:.2f}', {'size': 14}),
         (f'{nlow} of {n_poor} poor basins have <50% training Q', {'size': 14, 'bold': True, 'color': RED, 'space_before': 4}),
         ('Pearson r (NSE vs train cov) = +0.30', {'size': 13, 'color': GREY, 'space_before': 6}),
         ('Worst cases:', {'bold': True, 'size': 13, 'space_before': 10}),
         ('03572110 : 0.00  (no training Q)', {'size': 12.5, 'color': INK}),
         ('03559500 : 0.09', {'size': 12.5}),
         ('03400800 : 0.18', {'size': 12.5}),
         ('03593800 : 0.28', {'size': 12.5}),
         ]
textbox(s, Inches(9.1), Inches(1.25), Inches(4.0), Inches(5.0), lines)
callout(s, 'The model never saw enough labelled flow for these basins — '
           'long-lead skill is the first thing to collapse without training signal.')

# ---- Slide 3 : test-data scarcity (ruled out) ----------------------------
s = add_slide()
banner(s, 'Cause 2 — Test-data scarcity (ruled out)', 4)
pic(s, P_TEST, Inches(0.3), Inches(1.15), height=Inches(5.05))
med_pt = poor['test_cov'].median(); med_gd = good['test_cov'].median()
lines = [('Finding', {'bold': True, 'size': 16, 'color': GREEN_DK}),
         (f'Poor median test coverage = {med_pt:.2f}', {'size': 14, 'space_before': 6}),
         (f'Healthy median = {med_gd:.2f}', {'size': 14}),
         ('Essentially identical', {'size': 14, 'bold': True, 'color': RED, 'space_before': 4}),
         ('Pearson r (NSE vs test cov) = +0.07', {'size': 13, 'color': GREY, 'space_before': 6}),
         ('-> test-set size does NOT explain poor', {'size': 13, 'space_before': 4}),
         ('   long-lead NSE.', {'size': 13}),
         ('Only exception:', {'bold': True, 'size': 13, 'space_before': 10}),
         ('03572110 : test cov 0.39, n=168 obs', {'size': 12.5, 'color': RGBColor(*[int(x*255) for x in matplotlib.colors.to_rgb(C_BOTH)])}),
         ('  (also has zero training data)', {'size': 12, 'color': GREY}),
         ]
textbox(s, Inches(9.1), Inches(1.25), Inches(4.0), Inches(5.0), lines)
callout(s, 'Every poor basin (except 03572110) has a near-complete test record — '
           'the failure is upstream of evaluation, not a metric artefact.')

# ---- Slide 4 : hydrology / features --------------------------------------
s = add_slide()
banner(s, 'Cause 3 — Hydrology: small, flashy, low-baseflow', 5)
pic(s, P_FEATURES, Inches(0.25), Inches(1.15), height=Inches(5.0))
C_HYDRO_RGB = RGBColor(*[int(x*255) for x in matplotlib.colors.to_rgb(C_HYDRO)])
lines = [('Basins that fail WITH good data', {'bold': True, 'size': 15, 'color': GREEN_DK}),
         ('(pure hydrology, not data)', {'size': 12, 'color': GREY}),
         ('03602500  train .81 test 1.0  NSE .99->.24', {'size': 11.5, 'space_before': 6, 'color': C_HYDRO_RGB}),
         ('03466208  train .96 test .98  NSE .99->.43', {'size': 11.5, 'color': C_HYDRO_RGB}),
         ('03597590  train .99 test 1.0  NSE .98->.53', {'size': 11.5, 'color': C_HYDRO_RGB}),
         ('Poor vs healthy medians:', {'bold': True, 'size': 12.5, 'space_before': 8}),
         (f'area {poor["area_sqkm"].median():.0f} vs {good["area_sqkm"].median():.0f} km2   '
          f'BFI {poor["baseflow_index"].median():.0f} vs {good["baseflow_index"].median():.0f}   '
          f'slope {poor["slope_pct"].median():.1f} vs {good["slope_pct"].median():.1f}%', {'size': 12}),
         ('03602500 hydrograph below: sharp spikes', {'size': 11, 'color': GREY, 'space_before': 8}),
         ('decaying within hours — unresolvable 18 h out.', {'size': 11, 'color': GREY}),
         ]
textbox(s, Inches(8.55), Inches(1.25), Inches(4.6), Inches(3.0), lines)
# flashy time-series inset for the cleanest data-OK failure (right column, fits above callout)
ts = f'{TS_DIR}/timeseries_03602500_jan2019.png'
if os.path.exists(ts):
    pic(s, ts, Inches(8.55), Inches(4.35), width=Inches(4.55))
callout(s, 'Small + low-baseflow + steep-recession = flashy. With a full record the '
           '1 h skill is excellent, yet peaks vanish before an 18 h horizon arrives.',
        top=SH - Inches(0.95))

# ---- Slide 6 : synthesis table -------------------------------------------
s = add_slide()
banner(s, 'Synthesis — one cause per poor basin', 6)
rows = ['basin', 'repeat?', 'NSE 1h', 'NSE 18h', 'train cov', 'test cov', 'area km2', 'BFI', 'primary cause']
order = poor.sort_values(['cause', 'late_mean']).index
tbl_shape = s.shapes.add_table(len(order) + 1, len(rows),
                               Inches(0.4), Inches(1.15),
                               Inches(12.5), Inches(0.4) * (len(order) + 1))
tbl = tbl_shape.table
for j, h in enumerate(rows):
    c = tbl.cell(0, j); c.text = h
    c.fill.solid(); c.fill.fore_color.rgb = GREEN_DK
    pp = c.text_frame.paragraphs[0]; pp.font.size = Pt(12); pp.font.bold = True
    pp.font.color.rgb = WHITE; pp.alignment = PP_ALIGN.CENTER
for i, b in enumerate(order, start=1):
    r = m.loc[b]; cz = poor.loc[b, 'cause']
    vals = [b, 'Yes' if b in core else '—',
            f'{r["1h_ahead"]:.2f}', f'{r["18h_ahead"]:.2f}',
            f'{r["train_cov"]:.2f}', f'{r["test_cov"]:.2f}',
            f'{r["area_sqkm"]:.0f}', f'{r["baseflow_index"]:.0f}',
            CAUSE_LBL[cz]]
    fill = RGBColor(*[int(x*255) for x in matplotlib.colors.to_rgb(CAUSE_C[cz])])
    for j, v in enumerate(vals):
        c = tbl.cell(i, j); c.text = v
        pp = c.text_frame.paragraphs[0]; pp.font.size = Pt(11)
        pp.alignment = PP_ALIGN.CENTER
        if j == len(vals) - 1:
            c.fill.solid(); c.fill.fore_color.rgb = fill
            pp.font.color.rgb = WHITE; pp.font.bold = True
        if j == 1 and v == 'Yes':
            pp.font.color.rgb = RED; pp.font.bold = True
n_train = int((poor['cause'] == 'train').sum())
n_both  = int((poor['cause'] == 'both').sum())
n_hydro = int((poor['cause'] == 'hydro').sum())
callout(s, f'{n_train} training-scarce  ·  {n_both} train+test-scarce  ·  {n_hydro} '
           f'flashy-hydrology (data OK). Fixes differ: more labels vs. better '
           f'peak modelling / longer conditioning.', top=SH - Inches(1.15))

# ---- Slide 7 : improvement plan by root cause ----------------------------
s = add_slide()
banner(s, 'Improvement plan — fix by root cause', 7)
by_cause = {cz: [b for b in POOR if poor.loc[b, 'cause'] == cz]
            for cz in ['train', 'both', 'hydro']}
FIXES = {
    'train': ['Borrow signal from data-rich neighbours — regional / transfer learning, pretrain then fine-tune',
              'Targeted gauge in-fill; extend the training record',
              'Augment the few sparse-label basins'],
    'both':  ['Flag metrics as low-confidence (only ~168 test points)',
              'Collect new observations before trusting or deploying this basin',
              'Training-signal fixes also apply'],
    'hydro': ['Longer / richer conditioning window to capture fast rising limbs',
              'Peak-weighted or quantile loss',
              'Add a flashiness covariate, or a dedicated flashy-basin head'],
}
card_w = Inches(3.98); gap = Inches(0.30); left0 = Inches(0.40); top0 = Inches(1.25)
for i, cz in enumerate(['train', 'both', 'hydro']):
    left = Emu(int(left0) + i * (int(card_w) + int(gap)))
    col = RGBColor(*rgb255(CAUSE_C[cz]))
    # header
    hd = s.shapes.add_shape(5, left, top0, card_w, Inches(1.0))
    hd.fill.solid(); hd.fill.fore_color.rgb = col; hd.line.fill.background()
    tf = hd.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.text = CAUSE_LBL[cz]; p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(15); p.font.bold = True; p.font.color.rgb = WHITE
    p.font.name = 'Liberation Sans'
    p2 = tf.add_paragraph(); p2.text = f'{len(by_cause[cz])} basin' + ('s' if len(by_cause[cz]) != 1 else '')
    p2.alignment = PP_ALIGN.CENTER; p2.font.size = Pt(12); p2.font.color.rgb = WHITE
    # body
    body = [('Basins:  ' + ', '.join(by_cause[cz]), {'bold': True, 'size': 11.5, 'color': col})]
    body.append(('Fixes', {'bold': True, 'size': 12.5, 'color': col, 'space_before': 12}))
    for ln in FIXES[cz]:
        body.append(('•  ' + ln, {'size': 12, 'color': INK, 'space_before': 4}))
    textbox(s, Emu(int(left) + int(Inches(0.08))), Inches(2.5),
            Emu(int(card_w) - int(Inches(0.16))), Inches(4.0), body)
callout(s, 'Two failure families, two remedies: data-starved basins need more / borrowed '
           'labels; flashy basins need better fast-peak modelling. Test-set size is not a lever.',
        top=SH - Inches(1.0))

prs.save(PPTX_OUT)
print('saved', PPTX_OUT)
print('poor basin causes:')
print(poor[['late_mean', 'train_cov', 'test_cov', 'n_valid_test',
            'area_sqkm', 'baseflow_index', 'cause']].round(2).to_string())
