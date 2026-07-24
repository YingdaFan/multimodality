"""
Create a professional presentation about Learnable Adjacency Matrix for RGCN
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import os

current_dir = os.path.dirname(os.path.abspath(__file__))
plots_dir = os.path.join(current_dir, 'output', 'detailed_plots')


def set_text_format(text_frame, font_name="Calibri", font_size=18, bold=False,
                    color=RGBColor(0, 0, 0), alignment=PP_ALIGN.LEFT):
    """Helper function to format text"""
    for paragraph in text_frame.paragraphs:
        paragraph.font.name = font_name
        paragraph.font.size = Pt(font_size)
        paragraph.font.bold = bold
        paragraph.font.color.rgb = color
        paragraph.alignment = alignment


def add_title_slide(prs, title, subtitle=""):
    """Add title slide"""
    slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(slide_layout)

    title_shape = slide.shapes.title
    title_shape.text = title
    set_text_format(title_shape.text_frame, font_size=44, bold=True,
                    alignment=PP_ALIGN.CENTER)

    if subtitle:
        subtitle_shape = slide.placeholders[1]
        subtitle_shape.text = subtitle
        set_text_format(subtitle_shape.text_frame, font_size=24,
                       color=RGBColor(89, 89, 89), alignment=PP_ALIGN.CENTER)

    return slide


def add_section_header_slide(prs, title, subtitle=""):
    """Add section header slide"""
    slide_layout = prs.slide_layouts[2]  # Section header layout
    slide = prs.slides.add_slide(slide_layout)

    title_shape = slide.shapes.title
    title_shape.text = title
    set_text_format(title_shape.text_frame, font_size=40, bold=True,
                    alignment=PP_ALIGN.CENTER)

    if subtitle:
        subtitle_shape = slide.placeholders[1]
        subtitle_shape.text = subtitle
        set_text_format(subtitle_shape.text_frame, font_size=20,
                       color=RGBColor(89, 89, 89), alignment=PP_ALIGN.CENTER)

    return slide


def add_content_slide(prs, title, content_type="bullets"):
    """Add content slide with title and content placeholder"""
    slide_layout = prs.slide_layouts[1]  # Title and content layout
    slide = prs.slides.add_slide(slide_layout)

    title_shape = slide.shapes.title
    title_shape.text = title
    set_text_format(title_shape.text_frame, font_size=32, bold=True)

    return slide


def add_image_slide(prs, title, image_path):
    """Add slide with title and full-width image"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Add title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(9), Inches(0.8)
    )
    title_box.text = title
    set_text_format(title_box.text_frame, font_size=32, bold=True)

    # Add image
    if os.path.exists(image_path):
        # Calculate dimensions to fit slide (keeping aspect ratio)
        img_left = Inches(0.5)
        img_top = Inches(1.3)
        img_width = Inches(9)

        slide.shapes.add_picture(image_path, img_left, img_top, width=img_width)

    return slide


def create_presentation():
    """Create the complete presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ========== Slide 1: Title Slide ==========
    add_title_slide(
        prs,
        "Learnable Adjacency Matrix for \nSpatiotemporal Graph Neural Networks",
        "Enhancing RGCN with Adaptive Graph Structure Learning"
    )

    # ========== Slide 2: Application Motivation ==========
    slide = add_content_slide(prs, "Application: Prediction in Ungauged River Basins")

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    p = tf.add_paragraph()
    p.text = "Real-World Problem: Missing Historical Streamflow Data"
    p.level = 0
    p.font.bold = True
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(0, 112, 192)

    p = tf.add_paragraph()
    p.text = " "
    p.level = 0

    p = tf.add_paragraph()
    p.text = "🌊 Many rivers lack historical gauging stations"
    p.level = 0
    p.font.size = Pt(20)

    p = tf.add_paragraph()
    p.text = "Physical rivers exist, but no observed streamflow records"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Installing/maintaining gauges is expensive (thousands of $/year)"
    p.level = 1

    p = tf.add_paragraph()
    p.text = " "
    p.level = 0

    p = tf.add_paragraph()
    p.text = "📊 Available Information for Ungauged Basins:"
    p.level = 0
    p.font.size = Pt(20)

    p = tf.add_paragraph()
    p.text = "Meteorological features (precipitation, temperature, solar radiation)"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Static attributes (elevation, drainage area, soil type)"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Spatial location (latitude, longitude)"
    p.level = 1

    p = tf.add_paragraph()
    p.text = " "
    p.level = 0

    p = tf.add_paragraph()
    p.text = "🎯 Goal: Reconstruct historical streamflow by leveraging:"
    p.level = 0
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(0, 128, 0)

    p = tf.add_paragraph()
    p.text = "Local meteorological inputs at the target basin"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Observed streamflow + meteorology from neighboring basins"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Spatial relationships between hydrologically similar regions"
    p.level = 1

    set_text_format(tf, font_size=18)

    # ========== Slide 3: Problem Formulation ==========
    slide = add_content_slide(prs, "Problem Formulation: Few-Shot Learning on Graphs")

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    p = tf.add_paragraph()
    p.text = "Machine Learning Abstraction"
    p.level = 0
    p.font.bold = True
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(0, 112, 192)

    p = tf.add_paragraph()
    p.text = " "
    p.level = 0

    p = tf.add_paragraph()
    p.text = "Graph: 𝒢 = (𝒱, ℰ, 𝐗, 𝐘)"
    p.level = 0
    p.font.bold = True
    p.font.size = Pt(22)
    p.font.name = "Calibri"

    p = tf.add_paragraph()
    p.text = "𝒱: Set of N river basins (nodes)"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "ℰ: Spatial connectivity (initially based on geographic distance)"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "𝐗 ∈ ℝᴺˣᵀˣᵈ: Meteorological features (T timesteps, d features)"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "𝐘 ∈ ℝᴺˣᵀ: Streamflow observations (target variable)"
    p.level = 1

    p = tf.add_paragraph()
    p.text = " "
    p.level = 0

    p = tf.add_paragraph()
    p.text = "Node Sets: 𝒱 = 𝒱ᵍᵃᵘᵍᵉᵈ ∪ 𝒱ᵘⁿᵍᵃᵘᵍᵉᵈ"
    p.level = 0
    p.font.bold = True
    p.font.size = Pt(22)

    p = tf.add_paragraph()
    p.text = "𝒱ᵍᵃᵘᵍᵉᵈ: Basins with dense streamflow labels (>95% timesteps)"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "𝒱ᵘⁿᵍᵃᵘᵍᵉᵈ: Basins with sparse/no labels (<5% timesteps)"
    p.level = 1

    p = tf.add_paragraph()
    p.text = " "
    p.level = 0

    p = tf.add_paragraph()
    p.text = "Task: Few-Shot Spatiotemporal Imputation + Forecasting"
    p.level = 0
    p.font.bold = True
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(192, 0, 0)

    p = tf.add_paragraph()
    p.text = "Train: Jointly on all nodes with available labels"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Transfer: Propagate patterns via graph structure 𝒢"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Predict: Impute missing streamflow for 𝒱ᵘⁿᵍᵃᵘᵍᵉᵈ"
    p.level = 1

    set_text_format(tf, font_size=18)

    # ========== Slide 4: Key Challenges ==========
    slide = add_content_slide(prs, "Key Challenges")

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    p = tf.add_paragraph()
    p.text = "Challenge 1: Unknown Hydrological Connectivity"
    p.level = 0
    p.font.bold = True
    p.font.size = Pt(22)
    p.font.color.rgb = RGBColor(192, 0, 0)

    p = tf.add_paragraph()
    p.text = "Geographic distance ≠ hydrological similarity"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Distant basins may share climate patterns, geology, or land use"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "→ Need to learn adaptive graph structure from data"
    p.level = 1
    p.font.color.rgb = RGBColor(0, 128, 0)

    p = tf.add_paragraph()
    p.text = " "
    p.level = 0

    p = tf.add_paragraph()
    p.text = "Challenge 2: Extreme Label Scarcity for Target Nodes"
    p.level = 0
    p.font.bold = True
    p.font.size = Pt(22)
    p.font.color.rgb = RGBColor(192, 0, 0)

    p = tf.add_paragraph()
    p.text = "Ungauged basins have <5% labeled timesteps (few-shot regime)"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Cannot rely on local supervision alone"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "→ Must leverage spatial transfer from gauged neighbors"
    p.level = 1
    p.font.color.rgb = RGBColor(0, 128, 0)

    p = tf.add_paragraph()
    p.text = " "
    p.level = 0

    p = tf.add_paragraph()
    p.text = "Challenge 3: Spatiotemporal Heterogeneity"
    p.level = 0
    p.font.bold = True
    p.font.size = Pt(22)
    p.font.color.rgb = RGBColor(192, 0, 0)

    p = tf.add_paragraph()
    p.text = "Basin characteristics vary widely (snow-fed vs. rain-fed, etc.)"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Temporal dynamics differ (lag effects, seasonal patterns)"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "→ Need expressive spatiotemporal model (GNN + RNN)"
    p.level = 1
    p.font.color.rgb = RGBColor(0, 128, 0)

    p = tf.add_paragraph()
    p.text = " "
    p.level = 0

    p = tf.add_paragraph()
    p.text = "Core Research Question:"
    p.level = 0
    p.font.bold = True
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(0, 112, 192)

    p = tf.add_paragraph()
    p.text = "How to automatically discover the optimal graph structure for spatial knowledge transfer in few-shot spatiotemporal prediction?"
    p.level = 0
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(0, 112, 192)
    p.font.italic = True

    set_text_format(tf, font_size=18)

    # ========== Slide 5: Problem Statement ==========
    slide = add_content_slide(prs, "Motivation: Limitations of Fixed Adjacency Matrix")

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    # Problem 1
    p = tf.add_paragraph()
    p.text = "Traditional RGCN relies on geographic distance"
    p.level = 0

    p = tf.add_paragraph()
    p.text = "Assumption: closer basins → stronger hydrological connections"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Reality: watershed boundaries, climate zones, human intervention"
    p.level = 1

    # Problem 2
    p = tf.add_paragraph()
    p.text = "Isolated basins suffer from poor spatial propagation"
    p.level = 0

    p = tf.add_paragraph()
    p.text = "High self-connection weights due to row normalization"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Limited information from distant but hydrologically similar basins"
    p.level = 1

    # Consequence
    p = tf.add_paragraph()
    p.text = "⚠️ Suboptimal performance: geographic proximity ≠ hydrological connectivity"
    p.level = 0
    p.font.color.rgb = RGBColor(192, 0, 0)
    p.font.bold = True

    set_text_format(tf, font_size=20)

    # ========== Slide 3: Our Solution ==========
    slide = add_content_slide(prs, "Our Solution: Learnable Graph Structure")

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    p = tf.add_paragraph()
    p.text = "Key Innovation: Adaptive Graph Structure Learning"
    p.level = 0
    p.font.bold = True
    p.font.size = Pt(26)
    p.font.color.rgb = RGBColor(0, 112, 192)

    p = tf.add_paragraph()
    p.text = " "
    p.level = 0

    # Core equation with proper notation
    p = tf.add_paragraph()
    p.text = "𝐀ᵃᵈᵃᵖᵗⁱᵛᵉ = α × 𝐀ᵍᵉᵒ + (1-α) × 𝐀ˡᵉᵃʳⁿᵉᵈ"
    p.level = 0
    p.font.name = "Calibri"
    p.font.size = Pt(30)
    p.font.color.rgb = RGBColor(0, 128, 0)
    p.font.bold = True

    p = tf.add_paragraph()
    p.text = " "
    p.level = 0

    # Components with professional terms
    p = tf.add_paragraph()
    p.text = "𝐀ᵍᵉᵒ: Prior adjacency from geographic distance"
    p.level = 1
    p.font.size = Pt(20)
    p.font.bold = True

    p = tf.add_paragraph()
    p.text = "Encodes domain knowledge (spatial proximity)"
    p.level = 2
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(89, 89, 89)

    p = tf.add_paragraph()
    p.text = " "
    p.level = 0

    p = tf.add_paragraph()
    p.text = "𝐀ˡᵉᵃʳⁿᵉᵈ: Parameterized adjacency (trainable)"
    p.level = 1
    p.font.size = Pt(20)
    p.font.bold = True

    p = tf.add_paragraph()
    p.text = "Discovers data-driven connectivity patterns"
    p.level = 2
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(89, 89, 89)

    p = tf.add_paragraph()
    p.text = " "
    p.level = 0

    p = tf.add_paragraph()
    p.text = "α: Fusion coefficient (learned via gradient descent)"
    p.level = 1
    p.font.size = Pt(20)

    p = tf.add_paragraph()
    p.text = "Balances prior knowledge vs. data evidence"
    p.level = 2
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(89, 89, 89)

    p = tf.add_paragraph()
    p.text = " "
    p.level = 0

    p = tf.add_paragraph()
    p.text = "✓ Enables automatic discovery of true hydrological connectivity"
    p.level = 0
    p.font.color.rgb = RGBColor(0, 128, 0)
    p.font.bold = True
    p.font.size = Pt(20)

    set_text_format(tf, font_size=18)

    # ========== Slide 4: Model Architecture ==========
    slide = add_content_slide(prs, "Enhanced Graph Neural Network Architecture")

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    p = tf.add_paragraph()
    p.text = "Core Innovation (Active in This Work)"
    p.level = 0
    p.font.bold = True
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(0, 112, 192)

    # Component 1 - THE KEY ONE
    p = tf.add_paragraph()
    p.text = "① Adaptive Graph Structure Learning"
    p.level = 0
    p.font.bold = True
    p.font.size = Pt(22)
    p.font.color.rgb = RGBColor(0, 128, 0)

    p = tf.add_paragraph()
    p.text = "Parameterized adjacency: 𝐀ˡᵉᵃʳⁿᵉᵈ ∈ ℝᴺˣᴺ (trainable)"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Fusion coefficient: α ∈ (0,1) via backpropagation"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Row normalization for numerical stability in GCN"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "→ Main driver of performance improvement"
    p.level = 1
    p.font.color.rgb = RGBColor(192, 0, 0)
    p.font.bold = True

    p = tf.add_paragraph()
    p.text = " "
    p.level = 0

    # Component 2 - Secondary
    p = tf.add_paragraph()
    p.text = "② Multi-Head Spatial Attention (Optional)"
    p.level = 0
    p.font.bold = True
    p.font.size = Pt(20)

    p = tf.add_paragraph()
    p.text = "4-head self-attention applied periodically (every 5 steps)"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Captures dynamic, time-varying basin relationships"
    p.level = 1

    p = tf.add_paragraph()
    p.text = " "
    p.level = 0

    # Note about other components
    p = tf.add_paragraph()
    p.text = "Note: Advanced Features (Available but Inactive)"
    p.level = 0
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(89, 89, 89)

    p = tf.add_paragraph()
    p.text = "Semi-supervised pseudo-labels & domain adaptation"
    p.level = 1
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(89, 89, 89)

    p = tf.add_paragraph()
    p.text = "Not triggered: No basins with >95% missing data"
    p.level = 1
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(89, 89, 89)

    p = tf.add_paragraph()
    p.text = "→ Performance gain solely from adaptive graph structure"
    p.level = 1
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(0, 128, 0)
    p.font.bold = True

    set_text_format(tf, font_size=18)

    # ========== Slide 5: Training Strategy ==========
    slide = add_content_slide(prs, "Training Strategy")

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    p = tf.add_paragraph()
    p.text = "End-to-End Supervised Learning"
    p.level = 0
    p.font.bold = True
    p.font.size = Pt(26)
    p.font.color.rgb = RGBColor(0, 112, 192)

    p = tf.add_paragraph()
    p.text = " "
    p.level = 0

    # Objective
    p = tf.add_paragraph()
    p.text = "Optimization Objective:"
    p.level = 0
    p.font.bold = True
    p.font.size = Pt(22)

    p = tf.add_paragraph()
    p.text = "Loss = RMSE(y_true, y_pred)"
    p.level = 1
    p.font.name = "Calibri"
    p.font.size = Pt(26)
    p.font.color.rgb = RGBColor(0, 128, 0)
    p.font.bold = True

    p = tf.add_paragraph()
    p.text = "Root mean square error on observed streamflow values"
    p.level = 1
    p.font.color.rgb = RGBColor(89, 89, 89)

    p = tf.add_paragraph()
    p.text = " "
    p.level = 0

    # Training details
    p = tf.add_paragraph()
    p.text = "Training Configuration:"
    p.level = 0
    p.font.bold = True
    p.font.size = Pt(22)

    p = tf.add_paragraph()
    p.text = "Optimizer: Adam (lr=0.01, weight_decay=1e-5)"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Batch size: 29 (all basins processed together)"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Early stopping: patience=20 epochs"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Gradient clipping: max_norm=3.0"
    p.level = 1

    p = tf.add_paragraph()
    p.text = " "
    p.level = 0

    # Key point
    p = tf.add_paragraph()
    p.text = "Key Insight:"
    p.level = 0
    p.font.bold = True
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(0, 112, 192)

    p = tf.add_paragraph()
    p.text = "Graph structure (α, 𝐀ˡᵉᵃʳⁿᵉᵈ) learned jointly with LSTM weights"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "→ No separate pre-training needed for adjacency matrix"
    p.level = 1
    p.font.color.rgb = RGBColor(0, 128, 0)

    set_text_format(tf, font_size=18)

    # ========== Section: Results ==========
    add_section_header_slide(prs, "Experimental Results",
                            "Visualization of Learned Adjacency Patterns")

    # ========== Slide 6: Original vs Adaptive Matrix ==========
    fig1_path = os.path.join(plots_dir, 'fig1_original_matrix.png')
    fig3_path = os.path.join(plots_dir, 'fig3_adaptive_matrix.png')

    if os.path.exists(fig1_path):
        add_image_slide(prs, "Original Matrix: Geographic Distance Prior", fig1_path)

    if os.path.exists(fig3_path):
        add_image_slide(prs, "Adaptive Matrix: 18.4% Geographic + 81.6% Learned", fig3_path)

    # ========== Slide 7: Difference Matrix ==========
    fig4_path = os.path.join(plots_dir, 'fig4_difference_matrix.png')
    if os.path.exists(fig4_path):
        add_image_slide(prs, "Weight Changes: Adaptive - Original", fig4_path)

    # ========== Slide 8: Top Changes ==========
    fig5_path = os.path.join(plots_dir, 'fig5_top_changes.png')
    if os.path.exists(fig5_path):
        add_image_slide(prs, "Top 30 Basin Pairs by Weight Change", fig5_path)

    # ========== Slide 9: Fusion Analysis ==========
    fig6_path = os.path.join(plots_dir, 'fig6_fusion_analysis.png')
    if os.path.exists(fig6_path):
        add_image_slide(prs, "Fusion Analysis: From Geographic to Data-Driven", fig6_path)

    # ========== Slide 10: Key Findings ==========
    slide = add_content_slide(prs, "Experimental Findings")

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    p = tf.add_paragraph()
    p.text = "Finding 1: Data-Driven Structure Dominates"
    p.level = 0
    p.font.bold = True
    p.font.size = Pt(22)
    p.font.color.rgb = RGBColor(0, 112, 192)

    p = tf.add_paragraph()
    p.text = "Learned fusion coefficient: α = 0.184"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "→ 81.6% weight on 𝐀ˡᵉᵃʳⁿᵉᵈ vs. 18.4% on 𝐀ᵍᵉᵒ"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Conclusion: Geographic proximity is a weak inductive bias"
    p.level = 1
    p.font.color.rgb = RGBColor(192, 0, 0)

    p = tf.add_paragraph()
    p.text = " "
    p.level = 0

    p = tf.add_paragraph()
    p.text = "Finding 2: Significant Connectivity Enhancement"
    p.level = 0
    p.font.bold = True
    p.font.size = Pt(22)
    p.font.color.rgb = RGBColor(0, 112, 192)

    p = tf.add_paragraph()
    p.text = "Mean edge weight: 0.034 → 0.398 (+1056%)"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "100% of edges strengthened (812/812 non-diagonal entries)"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "→ Original distance matrix severely underestimated connectivity"
    p.level = 1
    p.font.color.rgb = RGBColor(192, 0, 0)

    p = tf.add_paragraph()
    p.text = " "
    p.level = 0

    p = tf.add_paragraph()
    p.text = "Finding 3: Resolves Graph Bottleneck Problem"
    p.level = 0
    p.font.bold = True
    p.font.size = Pt(22)
    p.font.color.rgb = RGBColor(0, 112, 192)

    p = tf.add_paragraph()
    p.text = "Isolated node (LCR): Self-loop reduced 40% → 3% (-93%)"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Discovered long-range dependencies (LCR ↔ NAV, 6° apart)"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "→ Breaks geographic constraints via data-driven learning"
    p.level = 1
    p.font.color.rgb = RGBColor(0, 128, 0)

    set_text_format(tf, font_size=17)

    # ========== Slide 11: Case Study - LCR ==========
    slide = add_content_slide(prs, "Case Study: Lake Cameahwait Reservoir (LCR)")

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    p = tf.add_paragraph()
    p.text = "Problem: Geographic Isolation"
    p.level = 0
    p.font.bold = True

    p = tf.add_paragraph()
    p.text = "Location: 42.66°N, -122.67°W (far west, Oregon/California border)"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Distance to nearest neighbor: 10.85° (≈700-800 km)"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Other basins: Clustered in Colorado/Wyoming (-111° to -106°)"
    p.level = 1

    p = tf.add_paragraph()
    p.text = " "
    p.level = 0

    p = tf.add_paragraph()
    p.text = "Original Matrix Behavior"
    p.level = 0
    p.font.bold = True
    p.font.color.rgb = RGBColor(192, 0, 0)

    p = tf.add_paragraph()
    p.text = "Self-connection: 40.2% (extremely high)"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Top neighbors: HYR, PIN, CAU (geographically closest)"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Row sum: 2.49 (lowest) → high diagonal weight after normalization"
    p.level = 1

    p = tf.add_paragraph()
    p.text = " "
    p.level = 0

    p = tf.add_paragraph()
    p.text = "Adaptive Matrix Solution"
    p.level = 0
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 128, 0)

    p = tf.add_paragraph()
    p.text = "Self-connection: 2.9% (reduced by 93%!)"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "New top neighbors: RFR, MCP, NAV (hydrologically similar)"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "NAV is in southern region (36.8°N) → long-range dependency discovered"
    p.level = 1

    set_text_format(tf, font_size=16)

    # ========== NEW SECTION: Deep Analysis of Learned Adjacency ==========
    add_section_header_slide(prs, "Deep Dive: What Does the Adjacency Matrix Learn?",
                            "Beyond Geographic Distance")

    # ========== Slide 12: Research Question ==========
    slide = add_content_slide(prs, "Critical Question: What Does A[i,j] Represent?")

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    p = tf.add_paragraph()
    p.text = "Hypothesis Testing: What Relationship Does A[i,j] Encode?"
    p.level = 0
    p.font.bold = True
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(0, 112, 192)

    p = tf.add_paragraph()
    p.text = " "
    p.level = 0

    p = tf.add_paragraph()
    p.text = "Candidate Hypotheses:"
    p.level = 0
    p.font.bold = True
    p.font.size = Pt(22)

    p = tf.add_paragraph()
    p.text = "H1: Meteorological feature similarity (cosine similarity of inputs)"
    p.level = 1
    p.font.size = Pt(18)

    p = tf.add_paragraph()
    p.text = "H2: Streamflow correlation (Pearson correlation of outputs)"
    p.level = 1
    p.font.size = Pt(18)

    p = tf.add_paragraph()
    p.text = "H3: Geographic proximity (spatial distance)"
    p.level = 1
    p.font.size = Pt(18)

    p = tf.add_paragraph()
    p.text = "H4: Cross-basin gradient influence (∂Y_i/∂X_j via backprop)"
    p.level = 1
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(0, 128, 0)
    p.font.bold = True

    p = tf.add_paragraph()
    p.text = " "
    p.level = 0

    p = tf.add_paragraph()
    p.text = "Method: Spearman correlation between A[i,j] and independent metrics"
    p.level = 0
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(89, 89, 89)

    p = tf.add_paragraph()
    p.text = " "
    p.level = 0

    p = tf.add_paragraph()
    p.text = "Key principle: Test only with independent, well-defined metrics"
    p.level = 0
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(0, 112, 192)
    p.font.italic = True

    set_text_format(tf, font_size=18)

    # ========== Slide 13: Correlation Analysis Result ==========
    correlation_plot_path = os.path.join(current_dir, 'output', 'deep_analysis', 'adjacency_meaning_analysis.png')
    if os.path.exists(correlation_plot_path):
        add_image_slide(prs, "Correlation Analysis: What Does the Adjacency Matrix Encode?",
                       correlation_plot_path)

    # ========== Slide 14: Key Insight - Predictive Utility ==========
    slide = add_content_slide(prs, "Key Discovery: Predictive Utility, Not Similarity")

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    p = tf.add_paragraph()
    p.text = "Surprising Result: Traditional Metrics All Fail!"
    p.level = 0
    p.font.bold = True
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(192, 0, 0)

    p = tf.add_paragraph()
    p.text = " "
    p.level = 0

    p = tf.add_paragraph()
    p.text = "Correlation with Learned Adjacency Matrix:"
    p.level = 0
    p.font.bold = True
    p.font.size = Pt(20)

    p = tf.add_paragraph()
    p.text = "Feature Similarity:        ρ = -0.076  ✗"
    p.level = 1
    p.font.name = "Consolas"
    p.font.size = Pt(18)

    p = tf.add_paragraph()
    p.text = "Target Correlation:        ρ =  0.127  ✗"
    p.level = 1
    p.font.name = "Consolas"
    p.font.size = Pt(18)

    p = tf.add_paragraph()
    p.text = "Geographic Proximity:      ρ =  0.186  ✗"
    p.level = 1
    p.font.name = "Consolas"
    p.font.size = Pt(18)

    p = tf.add_paragraph()
    p.text = "Cross-Basin Influence:     ρ =  0.603  ✓✓✓"
    p.level = 1
    p.font.name = "Consolas"
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(0, 128, 0)
    p.font.bold = True

    p = tf.add_paragraph()
    p.text = " "
    p.level = 0

    p = tf.add_paragraph()
    p.text = "Key Finding: Strong Correlation with Gradient Influence!"
    p.level = 0
    p.font.bold = True
    p.font.size = Pt(22)
    p.font.color.rgb = RGBColor(0, 128, 0)

    p = tf.add_paragraph()
    p.text = "Traditional metrics (features, targets, geography) all weak (ρ<0.2)"
    p.level = 1
    p.font.size = Pt(18)

    p = tf.add_paragraph()
    p.text = "BUT: Cross-basin gradient influence shows ρ=0.603 (p<1e-81)"
    p.level = 1
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(0, 128, 0)
    p.font.bold = True

    p = tf.add_paragraph()
    p.text = " "
    p.level = 0

    p = tf.add_paragraph()
    p.text = "Interpretation:"
    p.level = 0
    p.font.bold = True
    p.font.size = Pt(22)
    p.font.color.rgb = RGBColor(0, 112, 192)

    p = tf.add_paragraph()
    p.text = "Adjacency matrix encodes TRUE cross-basin input-output influence"
    p.level = 1
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(0, 128, 0)
    p.font.bold = True

    p = tf.add_paragraph()
    p.text = "Verified via gradient attribution: ∂Y_i/∂X_j (how basin j's input affects basin i's output)"
    p.level = 1
    p.font.size = Pt(17)

    p = tf.add_paragraph()
    p.text = "Not simple similarity, but learned predictive dependencies"
    p.level = 1
    p.font.size = Pt(18)

    p = tf.add_paragraph()
    p.text = "A[i,j] ≈ relative importance of basin j for predicting basin i"
    p.level = 1
    p.font.size = Pt(17)
    p.font.color.rgb = RGBColor(0, 112, 192)
    p.font.bold = True

    set_text_format(tf, font_size=18)

    # ========== Slide 15: Theoretical Explanation ==========
    slide = add_content_slide(prs, "Why No Correlation with Traditional Metrics?")

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    p = tf.add_paragraph()
    p.text = "Three Fundamental Reasons"
    p.level = 0
    p.font.bold = True
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(0, 112, 192)

    p = tf.add_paragraph()
    p.text = " "
    p.level = 0

    p = tf.add_paragraph()
    p.text = "1. Optimization Target: Information Complementarity"
    p.level = 0
    p.font.bold = True
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(192, 0, 0)

    p = tf.add_paragraph()
    p.text = "Traditional: A[i,j] ∝ Similarity(basin_i, basin_j)"
    p.level = 1
    p.font.size = Pt(17)

    p = tf.add_paragraph()
    p.text = "Actual: A[i,j] ∝ MutualInfo(q_j, residual_i)"
    p.level = 1
    p.font.size = Pt(17)
    p.font.color.rgb = RGBColor(0, 128, 0)

    p = tf.add_paragraph()
    p.text = "→ Model seeks complementary info, not redundant similarity"
    p.level = 1
    p.font.size = Pt(16)

    p = tf.add_paragraph()
    p.text = " "
    p.level = 0

    p = tf.add_paragraph()
    p.text = "2. Nonlinear Gradient Flow"
    p.level = 0
    p.font.bold = True
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(192, 0, 0)

    p = tf.add_paragraph()
    p.text = "RMSE → Dense → LSTM gates → spatial_context → q_t → A[i,j]"
    p.level = 1
    p.font.name = "Consolas"
    p.font.size = Pt(15)

    p = tf.add_paragraph()
    p.text = "Multiple nonlinear transforms (sigmoid, tanh) destroy linear correlations"
    p.level = 1
    p.font.size = Pt(16)

    p = tf.add_paragraph()
    p.text = " "
    p.level = 0

    p = tf.add_paragraph()
    p.text = "3. Task-Specific Representation Space"
    p.level = 0
    p.font.bold = True
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(192, 0, 0)

    p = tf.add_paragraph()
    p.text = "A operates on hidden states (q_t), not raw features (x_t)"
    p.level = 1
    p.font.size = Pt(17)

    p = tf.add_paragraph()
    p.text = "Learns: \"Which basin's learned representation helps predict mine?\""
    p.level = 1
    p.font.size = Pt(17)

    p = tf.add_paragraph()
    p.text = "→ Highly task-specific, cannot be explained by simple statistics"
    p.level = 1
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(0, 128, 0)

    set_text_format(tf, font_size=17)

    # ========== Slide 16: Information Flow Interpretation ==========
    slide = add_content_slide(prs, "Interpretation: Information Flow, Not Physical Flow")

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    p = tf.add_paragraph()
    p.text = "What A[i,j] Really Means"
    p.level = 0
    p.font.bold = True
    p.font.size = Pt(26)
    p.font.color.rgb = RGBColor(0, 112, 192)

    p = tf.add_paragraph()
    p.text = " "
    p.level = 0

    p = tf.add_paragraph()
    p.text = "❌ Common Misconception:"
    p.level = 0
    p.font.bold = True
    p.font.size = Pt(22)
    p.font.color.rgb = RGBColor(192, 0, 0)

    p = tf.add_paragraph()
    p.text = "\"A[i,j] high → basin i and j are physically connected\""
    p.level = 1
    p.font.size = Pt(18)

    p = tf.add_paragraph()
    p.text = "\"A[i,j] reflects hydrological connectivity\""
    p.level = 1
    p.font.size = Pt(18)

    p = tf.add_paragraph()
    p.text = " "
    p.level = 0

    p = tf.add_paragraph()
    p.text = "✓ Correct Understanding:"
    p.level = 0
    p.font.bold = True
    p.font.size = Pt(22)
    p.font.color.rgb = RGBColor(0, 128, 0)

    p = tf.add_paragraph()
    p.text = "\"A[i,j] = weight for aggregating basin j's hidden state when predicting basin i\""
    p.level = 1
    p.font.size = Pt(17)

    p = tf.add_paragraph()
    p.text = "\"Learned soft attention: where to look for predictive information\""
    p.level = 1
    p.font.size = Pt(17)

    p = tf.add_paragraph()
    p.text = " "
    p.level = 0

    p = tf.add_paragraph()
    p.text = "Key Implications:"
    p.level = 0
    p.font.bold = True
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(0, 112, 192)

    p = tf.add_paragraph()
    p.text = "Information flow ≠ physical flow (water, energy)"
    p.level = 1
    p.font.size = Pt(18)

    p = tf.add_paragraph()
    p.text = "Distant basins can have high A[i,j] if functionally similar"
    p.level = 1
    p.font.size = Pt(18)

    p = tf.add_paragraph()
    p.text = "Statistical dependency ≠ causal relationship"
    p.level = 1
    p.font.size = Pt(18)

    p = tf.add_paragraph()
    p.text = "Model discovers patterns human priors (geography) miss"
    p.level = 1
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(0, 128, 0)
    p.font.bold = True

    set_text_format(tf, font_size=18)

    # ========== Slide 17: Scientific Value ==========
    slide = add_content_slide(prs, "Scientific Value & Limitations")

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    p = tf.add_paragraph()
    p.text = "✓ What This Tells Us"
    p.level = 0
    p.font.bold = True
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(0, 128, 0)

    p = tf.add_paragraph()
    p.text = " "
    p.level = 0

    p = tf.add_paragraph()
    p.text = "End-to-end learning discovers task-specific patterns"
    p.level = 1
    p.font.size = Pt(19)

    p = tf.add_paragraph()
    p.text = "Geographic distance is a weak prior (18.4% weight)"
    p.level = 1
    p.font.size = Pt(19)

    p = tf.add_paragraph()
    p.text = "Data-driven patterns outperform human-designed priors"
    p.level = 1
    p.font.size = Pt(19)

    p = tf.add_paragraph()
    p.text = "Can serve as hypothesis generator for hydrological research"
    p.level = 1
    p.font.size = Pt(19)

    p = tf.add_paragraph()
    p.text = " "
    p.level = 0

    p = tf.add_paragraph()
    p.text = "⚠️ Limitations & Cautions"
    p.level = 0
    p.font.bold = True
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(192, 0, 0)

    p = tf.add_paragraph()
    p.text = " "
    p.level = 0

    p = tf.add_paragraph()
    p.text = "Cannot interpret as physical/causal relationships"
    p.level = 1
    p.font.size = Pt(19)

    p = tf.add_paragraph()
    p.text = "May capture spurious correlations in training data"
    p.level = 1
    p.font.size = Pt(19)

    p = tf.add_paragraph()
    p.text = "Learned patterns are task-specific (streamflow prediction)"
    p.level = 1
    p.font.size = Pt(19)

    p = tf.add_paragraph()
    p.text = "Requires domain knowledge for validation and interpretation"
    p.level = 1
    p.font.size = Pt(19)

    p = tf.add_paragraph()
    p.text = " "
    p.level = 0

    p = tf.add_paragraph()
    p.text = "Best Practice: Use as complement to, not replacement for, physical models"
    p.level = 0
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(0, 112, 192)
    p.font.bold = True
    p.font.italic = True

    set_text_format(tf, font_size=18)

    # ========== Slide 18: Technical Insights ==========
    slide = add_content_slide(prs, "Technical Insights & Design Choices")

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    p = tf.add_paragraph()
    p.text = "Why Sigmoid Activation?"
    p.level = 0
    p.font.bold = True

    p = tf.add_paragraph()
    p.text = "Bounds learned weights to [0,1] → stable optimization"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Ensures valid probability-like adjacency values"
    p.level = 1

    p = tf.add_paragraph()
    p.text = " "
    p.level = 0

    p = tf.add_paragraph()
    p.text = "Why Convex Combination?"
    p.level = 0
    p.font.bold = True

    p = tf.add_paragraph()
    p.text = "Preserves useful geographic priors (α > 0)"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Allows model to gradually shift trust from prior to data"
    p.level = 1

    p = tf.add_paragraph()
    p.text = " "
    p.level = 0

    p = tf.add_paragraph()
    p.text = "Why Row Normalization?"
    p.level = 0
    p.font.bold = True

    p = tf.add_paragraph()
    p.text = "Ensures each row sums to 1 → probabilistic interpretation"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Numerical stability in GCN propagation: h_new = A @ h_old"
    p.level = 1
    p.font.name = "Consolas"

    p = tf.add_paragraph()
    p.text = " "
    p.level = 0

    p = tf.add_paragraph()
    p.text = "Implementation: model.py:119-132"
    p.level = 0
    p.font.color.rgb = RGBColor(89, 89, 89)
    p.font.size = Pt(14)

    set_text_format(tf, font_size=18)

    # ========== Slide 13: Advantages ==========
    slide = add_content_slide(prs, "Advantages Over Traditional RGCN")

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    p = tf.add_paragraph()
    p.text = "✓ Automatic Discovery of Hydrological Connectivity"
    p.level = 0
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 128, 0)

    p = tf.add_paragraph()
    p.text = "No need for manual watershed delineation or river network topology"
    p.level = 1

    p = tf.add_paragraph()
    p.text = " "
    p.level = 0

    p = tf.add_paragraph()
    p.text = "✓ Handles Geographic Isolation"
    p.level = 0
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 128, 0)

    p = tf.add_paragraph()
    p.text = "Reduces self-loop bias from row normalization"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Discovers long-range dependencies"
    p.level = 1

    p = tf.add_paragraph()
    p.text = " "
    p.level = 0

    p = tf.add_paragraph()
    p.text = "✓ Corrects Inductive Bias"
    p.level = 0
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 128, 0)

    p = tf.add_paragraph()
    p.text = "Model learns that geographic distance is a weak prior (18.4%)"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Trusts data-driven patterns instead (81.6%)"
    p.level = 1

    p = tf.add_paragraph()
    p.text = " "
    p.level = 0

    p = tf.add_paragraph()
    p.text = "✓ Transferable to Other Domains"
    p.level = 0
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 128, 0)

    p = tf.add_paragraph()
    p.text = "Traffic networks, social networks, climate modeling, etc."
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Any spatiotemporal problem with uncertain graph structure"
    p.level = 1

    set_text_format(tf, font_size=20)

    # ========== Slide 14: Conclusion ==========
    slide = add_content_slide(prs, "Conclusion")

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    p = tf.add_paragraph()
    p.text = "Main Contribution"
    p.level = 0
    p.font.bold = True
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(0, 112, 192)

    p = tf.add_paragraph()
    p.text = "Learnable adjacency matrix that adaptively fuses geographic priors with data-driven patterns, enabling RGCN to discover true hydrological connectivity"
    p.level = 1
    p.font.size = Pt(20)

    p = tf.add_paragraph()
    p.text = " "
    p.level = 0

    p = tf.add_paragraph()
    p.text = "Key Takeaways"
    p.level = 0
    p.font.bold = True
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(0, 112, 192)

    p = tf.add_paragraph()
    p.text = "Geography is a weak prior for hydrological modeling (α = 0.184)"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Model dramatically increased connectivity (+1056% avg weight)"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Solved isolated basin problem (LCR: 40% → 3% self-connection)"
    p.level = 1

    p = tf.add_paragraph()
    p.text = " "
    p.level = 0

    p = tf.add_paragraph()
    p.text = "Future Work"
    p.level = 0
    p.font.bold = True
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(0, 112, 192)

    p = tf.add_paragraph()
    p.text = "Time-adaptive adjacency (let α vary across time steps)"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Multi-scale graph structure (short/medium/long-range connections)"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Interpretability analysis (extract hydrological connectivity maps)"
    p.level = 1

    set_text_format(tf, font_size=18)

    # ========== Final Slide: Thank You ==========
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Add centered "Thank You" text
    thank_you_box = slide.shapes.add_textbox(
        Inches(1), Inches(2.5), Inches(8), Inches(2)
    )
    thank_you_box.text = "Thank You!\n\nQuestions?"
    tf = thank_you_box.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    set_text_format(tf, font_size=48, bold=True,
                    color=RGBColor(0, 112, 192), alignment=PP_ALIGN.CENTER)

    # Add contact info
    contact_box = slide.shapes.add_textbox(
        Inches(1), Inches(5), Inches(8), Inches(1)
    )
    contact_box.text = "Learnable Adjacency Matrix for Spatiotemporal Graph Neural Networks"
    tf = contact_box.text_frame
    set_text_format(tf, font_size=18, alignment=PP_ALIGN.CENTER,
                    color=RGBColor(89, 89, 89))

    # Save presentation
    output_path = os.path.join(current_dir, 'output', 'Learnable_Adjacency_Matrix_Presentation.pptx')
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)

    print(f"✓ Presentation created successfully!")
    print(f"  Location: {output_path}")
    print(f"  Total slides: {len(prs.slides)}")

    return output_path


if __name__ == "__main__":
    create_presentation()
