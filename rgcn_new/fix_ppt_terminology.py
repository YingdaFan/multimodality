#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
Fix terminology in the presentation:
Replace 'ungauged' with more accurate semi-supervised learning terminology
"""

from pptx import Presentation
from pptx.util import Pt
import re

def replace_text_in_runs(paragraph, old_text, new_text):
    """Replace text in paragraph while preserving formatting"""
    # Get full text
    full_text = ''.join(run.text for run in paragraph.runs)

    # Case-insensitive search
    pattern = re.compile(re.escape(old_text), re.IGNORECASE)
    if not pattern.search(full_text):
        return False

    # Replace
    new_full_text = pattern.sub(new_text, full_text)

    # Clear all runs except first
    for run in paragraph.runs[1:]:
        run.text = ""

    # Set new text to first run
    if paragraph.runs:
        paragraph.runs[0].text = new_full_text
        return True

    return False


def process_shape(shape, replacements):
    """Process a shape and replace text"""
    if not hasattr(shape, "text_frame"):
        return 0

    changes = 0
    for paragraph in shape.text_frame.paragraphs:
        for old_text, new_text in replacements.items():
            if replace_text_in_runs(paragraph, old_text, new_text):
                changes += 1
                print(f"    ✓ '{old_text}' → '{new_text}'")

    return changes


def main():
    ppt_path = "output/Learnable_Adjacency_Matrix_Presentation_Updated.pptx"
    prs = Presentation(ppt_path)

    print(f"Loaded presentation with {len(prs.slides)} slides\n")
    print("="*70)

    # Define replacements - order matters!
    replacements = {
        # Core terminology
        "Prediction in Ungauged River Basins": "Semi-Supervised Prediction for River Basins with Missing Labels",
        "Ungauged River Basins": "River Basins with Missing Labels",
        "ungauged basins": "unlabeled basins",
        "ungauged basin": "unlabeled basin",
        "Ungauged": "Unlabeled",

        # Gauged terminology
        "gauged basins": "labeled basins",
        "gauged basin": "labeled basin",
        "gauging stations": "monitoring stations with labels",
        "gauging station": "monitoring station with labels",
        "lack historical gauging": "lack historical observations",

        # Data terminology
        "Missing Historical Streamflow Data": "Missing Streamflow Labels",
        "missing data": "missing labels",

        # Learning paradigm
        "Transfer learning across basins": "Semi-supervised learning across basins",
        "Transfer Learning": "Semi-Supervised Learning",
        "transfer learning": "semi-supervised learning",
    }

    total_changes = 0
    slides_modified = 0

    for i, slide in enumerate(prs.slides, 1):
        print(f"\n📄 Slide {i}:")
        slide_changes = 0

        for shape in slide.shapes:
            shape_changes = process_shape(shape, replacements)
            slide_changes += shape_changes
            total_changes += shape_changes

        if slide_changes > 0:
            slides_modified += 1
            print(f"  ✅ {slide_changes} replacement(s) made")
        else:
            print(f"  ⏭️  No changes needed")

    # Save updated presentation
    output_path = "output/Learnable_Adjacency_Matrix_Presentation_Fixed.pptx"
    prs.save(output_path)

    print("\n" + "="*70)
    print(f"✅ Terminology correction completed!")
    print(f"   - Total replacements: {total_changes}")
    print(f"   - Slides modified: {slides_modified}/{len(prs.slides)}")
    print(f"   - Saved to: {output_path}")
    print("="*70)

    # Print summary
    print("\n📋 Terminology Changes Summary:")
    print("-"*70)
    key_changes = [
        ("Ungauged basins", "Unlabeled basins"),
        ("Gauged basins", "Labeled basins"),
        ("Missing data", "Missing labels"),
        ("Transfer learning", "Semi-supervised learning"),
    ]
    for old, new in key_changes:
        print(f"  • '{old}' → '{new}'")

    print("\n🎯 Key Conceptual Changes:")
    print("-"*70)
    print("  ✓ Task: Ungauged prediction → Semi-supervised prediction")
    print("  ✓ Basin types: Gauged/Ungauged → Labeled/Unlabeled")
    print("  ✓ Problem: Missing data → Missing labels")
    print("  ✓ Method: Transfer learning → Semi-supervised learning")
    print("  ✓ Focus: Emphasize that targets have features but no labels")


if __name__ == "__main__":
    main()
