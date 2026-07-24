#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
Add slides to Learnable_Adjacency_Matrix_Presentation.pptx
Based on insights from ATTENTION_MECHANISM_ANALYSIS.md
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

def add_title_and_content_with_table(prs, title_text, content_before_table, table_data, content_after_table=None):
    """Add a slide with title, text, table, and optional text after table"""

    # Use blank layout (6) for more control
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title_text
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 51, 102)

    # Add content before table
    top_position = Inches(1.2)
    for item in content_before_table:
        text_box = slide.shapes.add_textbox(Inches(0.5), top_position, Inches(9), Inches(0.4))
        text_frame = text_box.text_frame
        text_frame.word_wrap = True

        p = text_frame.paragraphs[0]
        if item.startswith('**') and item.endswith('**'):
            # Bold text
            p.text = item.strip('**').strip(':')
            p.font.bold = True
            p.font.size = Pt(18)
        elif item.startswith('- **'):
            # Bullet with bold part
            parts = item[4:].split('**:', 1)
            p.text = "• " + parts[0] + ": " + parts[1] if len(parts) > 1 else "• " + item[4:].strip('**')
            p.font.size = Pt(16)
        else:
            p.text = item
            p.font.size = Pt(16)

        top_position += Inches(0.35)

    # Add table
    top_position += Inches(0.1)
    rows, cols = len(table_data), len(table_data[0])
    left = Inches(0.5)
    width = Inches(9)
    height = Inches(rows * 0.4)

    table_shape = slide.shapes.add_table(rows, cols, left, top_position, width, height)
    table = table_shape.table

    # Fill table
    for i, row_data in enumerate(table_data):
        for j, cell_text in enumerate(row_data):
            cell = table.rows[i].cells[j]
            cell.text = str(cell_text)

            # Format text
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(12)
                paragraph.alignment = PP_ALIGN.CENTER

                # Header row
                if i == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = RGBColor(255, 255, 255)
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(0, 51, 102)

    top_position += height + Inches(0.15)

    # Add content after table
    if content_after_table:
        for item in content_after_table:
            text_box = slide.shapes.add_textbox(Inches(0.5), top_position, Inches(9), Inches(0.35))
            text_frame = text_box.text_frame
            text_frame.word_wrap = True

            p = text_frame.paragraphs[0]
            if item.startswith('**'):
                # Bold text
                run = p.add_run()
                run.text = item.strip('**')
                run.font.bold = True
                run.font.size = Pt(16)
            elif item.startswith('- ✓'):
                p.text = item
                p.font.size = Pt(14)
            else:
                p.text = item
                p.font.size = Pt(14)

            top_position += Inches(0.3)

    return slide


def main():
    # Load existing presentation
    ppt_path = "./output/Learnable_Adjacency_Matrix_Presentation.pptx"
    prs = Presentation(ppt_path)

    print(f"Loaded presentation with {len(prs.slides)} slides")

    # ========== Slide 1: Core Discovery ==========
    title1 = "核心发现：可学习邻接矩阵的本质"

    content_before_table1 = [
        "**主要结论**: 可学习邻接矩阵本质上是一种静态Attention机制",
        "**α混合对应**: Prior-Guided Attention (文献中的经典技巧)",
        "",
        "**对比分析:**"
    ]

    table_data1 = [
        ["特性", "可学习邻接矩阵", "标准Attention"],
        ["权重计算", "直接参数化学习", "动态计算 (Q·K^T)"],
        ["加权聚合", "A @ q_t", "attn_weights @ V"],
        ["优化目标", "任务损失端到端优化", "任务损失端到端优化"],
        ["适用场景", "静态图结构学习", "动态序列建模"]
    ]

    content_after_table1 = [
        "",
        "**核心相似性:**",
        "- ✓ 都是加权聚合机制",
        "- ✓ 都学习'应该关注哪里'",
        "- ✓ 都可解释为软注意力"
    ]

    add_title_and_content_with_table(prs, title1, content_before_table1, table_data1, content_after_table1)
    print("Added Slide 1: Core Discovery")

    # ========== Slide 2: α Interpretation ==========
    title2 = "α ≈ 0.18 的深层含义"

    content_before_table2 = [
        "**多重理论解释:**"
    ]

    table_data2 = [
        ["视角", "解释", "α≈0.18的含义"],
        ["Bayesian", "Prior weight", "先验信念占后验的18%"],
        ["信息论", "Information ratio", "地理先验信息量是数据驱动的1/4.5"],
        ["优化理论", "Bias-Variance Tradeoff", "最优权衡点"],
        ["物理意义", "Prior trust", "地理距离可信度18%"]
    ]

    content_after_table2 = [
        "",
        "**关键发现:**",
        "• 不是失败，而是科学发现: 地理距离是弱先验的定量证据",
        "• 数据驱动模式提供了 4.4倍 于地理先验的信息量",
        "",
        "**凸组合的必要性:**",
        "• 单独 A_geo: 高偏差（欠拟合）",
        "• 单独 A_learned: 高方差（优化不稳定）",
        "• α混合: 最优权衡 + 物理可解释性"
    ]

    add_title_and_content_with_table(prs, title2, content_before_table2, table_data2, content_after_table2)
    print("Added Slide 2: α Interpretation")

    # Save the updated presentation
    output_path = "./output/Learnable_Adjacency_Matrix_Presentation_Updated.pptx"
    prs.save(output_path)

    print(f"\nPresentation updated successfully!")
    print(f"Total slides: {len(prs.slides)}")
    print(f"Saved to: {output_path}")
    print("\nYou can also overwrite the original file if needed.")


if __name__ == "__main__":
    main()
