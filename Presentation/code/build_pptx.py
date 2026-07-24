"""
Build the SAFER-Hydro / TRB deck as a single .pptx.

Slide 1:  cover (rendered page 1 of InflowForecast.pdf, with the
          collaborator line rewritten to put Yingda Fan first).
Slides 2-11: the ten polished figures, in the agreed order, each
          centred and scaled to fill a 16:9 slide.

Output: Presentation/presentation/SAFER_Hydro_TRB.pptx
"""

import os
from pathlib import Path

import fitz                                   # PyMuPDF
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches, Emu, Pt
from pptx.enum.shapes import MSO_SHAPE


# script lives in Presentation/code/ — go two up to reach imputation/
PROJECT_DIR = Path(__file__).resolve().parents[2]
PDF_PATH    = PROJECT_DIR / 'InflowForecast.pdf'
OUT_PPTX    = PROJECT_DIR / 'Presentation' / 'presentation' / 'SAFER_Hydro_TRB.pptx'
COVER_PNG   = PROJECT_DIR / 'Presentation' / 'presentation' / '_cover.png'


# Order of content slides (slide 2 onward)
CONTENT_PNGS = [
    PROJECT_DIR / 'Presentation' / 'presentation' / 'global_vs_tn_map.png',
    PROJECT_DIR / 'Presentation' / 'presentation' / 'data_atlas.png',
    PROJECT_DIR / 'Presentation' / 'presentation' / 'experiment_setting.png',
    PROJECT_DIR / 'Presentation' / 'presentation' / 'model_card.png',
    PROJECT_DIR / 'Presentation' / 'presentation' / 'coverage_vs_nse.png',
    PROJECT_DIR / 'Presentation' / 'presentation' / 'basin_panel.png',
    PROJECT_DIR / 'Presentation' / 'tennessee_18lead' / 'figures' / 'forecast_boxplot.png',
    PROJECT_DIR / 'Presentation' / 'tennessee_18lead' / 'figures' / 'forecast_cdf.png',
    PROJECT_DIR / 'Presentation' / 'tennessee_18lead' / 'figures' / 'forecast_lines.png',
    PROJECT_DIR / 'Presentation' / 'tennessee_18lead' / 'figures' / 'page9_summary.png',
]


# ============================== 1. Build the cover image ==============================
def build_cover_png():
    """Render page 1 of the PDF, then paint over the collaborator line and
    rewrite it with Yingda Fan first."""
    doc  = fitz.open(PDF_PATH)
    page = doc[0]
    src_w, src_h = page.rect.width, page.rect.height           # PDF points
    # Render at ~200 dpi for a crisp slide background
    mat  = fitz.Matrix(200 / 72, 200 / 72)
    pix  = page.get_pixmap(matrix=mat, alpha=False)
    img  = Image.frombytes('RGB', (pix.width, pix.height), pix.samples)
    doc.close()
    W, H = img.size
    print(f'cover render: {W} × {H} px')
    sx = W / src_w
    sy = H / src_h

    # Bounding box in PDF points, located via fitz text-search on this same
    # page (page is 959.76 × 540 points, line at y ≈ 448-471)
    BBOX_PDF = (33, 446, 470, 475)
    x0, y0, x1, y1 = (int(BBOX_PDF[0] * sx), int(BBOX_PDF[1] * sy),
                      int(BBOX_PDF[2] * sx), int(BBOX_PDF[3] * sy))

    # White-out the original line
    draw = ImageDraw.Draw(img)
    draw.rectangle((x0, y0, x1, y1), fill='white')

    # Re-draw with Yingda Fan first.  Use a sans-serif font close to the rest
    # of the deck — Liberation Sans is available on this box.
    font_size = int(0.55 * (y1 - y0))   # match the deleted line height
    font_paths = [
        '/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf',
        '/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf',
    ]
    font = None
    for fp in font_paths:
        if os.path.exists(fp):
            font = ImageFont.truetype(fp, font_size)
            break
    if font is None:
        font = ImageFont.load_default()

    new_text = 'Yingda Fan; Vinh Tran; Soumendra Bhanja'
    # vertical centring inside the box
    bbox  = draw.textbbox((0, 0), new_text, font=font)
    th    = bbox[3] - bbox[1]
    ty    = y0 + ((y1 - y0) - th) // 2 - bbox[1]
    draw.text((x0, ty), new_text, font=font, fill='black')

    COVER_PNG.parent.mkdir(parents=True, exist_ok=True)
    img.save(COVER_PNG, 'PNG')
    print(f'cover saved → {COVER_PNG}')


# ============================== 2. Assemble the .pptx ==============================
def build_pptx():
    prs = Presentation()
    # 16:9 slide dimensions
    prs.slide_width  = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank = prs.slide_layouts[6]
    sw, sh = prs.slide_width, prs.slide_height

    def add_image_slide(png_path, fit='cover'):
        """`cover`: scale the PNG to cover the whole slide, possibly
        cropping; `contain`: scale to fit inside the slide with letterboxing.
        Cover is right for the cover image; contain is right for figures we
        want to keep entirely visible."""
        slide = prs.slides.add_slide(blank)
        with Image.open(png_path) as im:
            iw, ih = im.size
        s_ratio = sw / sh
        i_ratio = iw / ih

        if fit == 'cover':
            if i_ratio > s_ratio:
                h_new = sh
                w_new = int(h_new * i_ratio)
            else:
                w_new = sw
                h_new = int(w_new / i_ratio)
        else:   # contain
            if i_ratio > s_ratio:
                w_new = sw
                h_new = int(w_new / i_ratio)
            else:
                h_new = sh
                w_new = int(h_new * i_ratio)
        x = (sw - w_new) // 2
        y = (sh - h_new) // 2
        slide.shapes.add_picture(str(png_path), x, y, w_new, h_new)
        return slide

    # Slide 1: cover (cover fit so it fills the slide)
    add_image_slide(COVER_PNG, fit='cover')

    # Slides 2-11: content figures (contain fit so nothing is cropped)
    for png in CONTENT_PNGS:
        if not png.exists():
            raise FileNotFoundError(png)
        add_image_slide(png, fit='contain')

    OUT_PPTX.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_PPTX)
    print(f'pptx saved → {OUT_PPTX}   ({len(prs.slides)} slides)')


if __name__ == '__main__':
    build_cover_png()
    build_pptx()
